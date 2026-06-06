"""Build a customer-facing PPT for data-standard lifecycle management."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "client_deliverables" / (
    "数据标准全生命周期智能化管理_客户讲解版.pptx"
)

W, H = Inches(13.333), Inches(7.5)
FONT = "Microsoft YaHei"

COLORS = {
    "bg": RGBColor(247, 250, 249),
    "ink": RGBColor(26, 35, 45),
    "muted": RGBColor(92, 105, 118),
    "line": RGBColor(218, 227, 231),
    "green": RGBColor(16, 118, 91),
    "green2": RGBColor(31, 143, 105),
    "blue": RGBColor(35, 83, 134),
    "amber": RGBColor(180, 119, 44),
    "red": RGBColor(153, 71, 68),
    "white": RGBColor(255, 255, 255),
    "pale_green": RGBColor(232, 246, 239),
    "pale_blue": RGBColor(232, 240, 249),
    "pale_amber": RGBColor(250, 243, 230),
}


def c(name: str) -> RGBColor:
    return COLORS[name]


def apply_run(run, size=16, color="ink", bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = c(color)
    run.font.bold = bold


def fill(shape, color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = c(color)


def border(shape, color="line", width=0.8):
    shape.line.color.rgb = c(color)
    shape.line.width = Pt(width)


def bg(slide):
    fill(slide.background, "bg")


def text(slide, x, y, w, h, value, size=16, color="ink",
         bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    apply_run(r, size=size, color=color, bold=bold)
    return box


def bullets(slide, x, y, w, h, items, size=15, color="ink", gap=7):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.text = ""
        r1 = p.add_run()
        r1.text = "• "
        apply_run(r1, size=size, color="green", bold=True)
        r2 = p.add_run()
        r2.text = item
        apply_run(r2, size=size, color=color)
    return box


def shape(slide, x, y, w, h, color="white", line="line", rounded=True):
    kind = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    s = slide.shapes.add_shape(kind, x, y, w, h)
    fill(s, color)
    border(s, line)
    return s


def header(slide, title: str, page: int):
    text(slide, Inches(0.55), Inches(0.22), Inches(9.6), Inches(0.42),
         title, size=18, color="green", bold=True)
    text(slide, Inches(11.3), Inches(0.24), Inches(1.8), Inches(0.25),
         "客户讲解版", size=10, color="muted", align=PP_ALIGN.RIGHT)
    conn = slide.shapes.add_connector(
        1, Inches(0.55), Inches(0.72), Inches(12.75), Inches(0.72)
    )
    conn.line.color.rgb = c("line")
    text(slide, Inches(0.55), Inches(7.08), Inches(7.0), Inches(0.22),
         "GIS Data Agent · 数据标准全生命周期智能化管理", size=9, color="muted")
    text(slide, Inches(12.15), Inches(7.08), Inches(0.9), Inches(0.22),
         f"{page:02d}", size=9, color="muted", align=PP_ALIGN.RIGHT)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    shape(slide, Inches(0), Inches(0), W, H, color="bg", line="bg",
          rounded=False)
    shape(slide, Inches(0.72), Inches(0.75), Inches(3.1), Inches(0.08),
          color="green", line="green", rounded=False)
    text(slide, Inches(0.72), Inches(1.35), Inches(10.8), Inches(0.75),
         "数据标准全生命周期智能化管理", size=34, color="ink", bold=True)
    text(slide, Inches(0.75), Inches(2.25), Inches(10.5), Inches(0.45),
         "把标准从静态文档变成可治理、可追溯、可驱动下游系统的数字资产",
         size=20, color="green", bold=True)
    bullets(slide, Inches(0.82), Inches(3.18), Inches(8.9), Inches(1.5), [
        "面向自然资源、测绘地理信息、国土空间规划等标准密集型业务",
        "覆盖采集、结构化、起草、审定、发布、派生、变更治理和共享复用",
        "让标准成为语义层、数据质检、缺陷分类、数据建模的单一权威源",
    ], size=18)
    for i, (title, body, color) in enumerate([
        ("标准可用", "条款、术语、数据元、值域结构化", "green"),
        ("治理可控", "审定、版本、影响分析、回滚", "blue"),
        ("下游可驱动", "语义层、质检规则、数据模型", "amber"),
    ]):
        x = Inches(0.82 + i * 3.55)
        shape(slide, x, Inches(5.38), Inches(3.0), Inches(0.95),
              color="white", line="line")
        text(slide, x + Inches(0.18), Inches(5.53), Inches(2.45),
             Inches(0.28), title, size=16, color=color, bold=True)
        text(slide, x + Inches(0.18), Inches(5.92), Inches(2.55),
             Inches(0.28), body, size=11, color="muted")


def pain_goal(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "为什么需要这套能力", page)
    columns = [
        ("传统痛点", [
            "标准散落在 Word、PDF、XMI 和人工经验中",
            "字段命名、值域解释、质检规则靠专家手工落地",
            "标准升版后，下游影响面难以排查",
            "智能问数/NL2SQL 缺少稳定的业务语义依据",
        ], "red", "pale_amber"),
        ("建设目标", [
            "把标准结构化为可计算的条款、术语、数据元和值域",
            "把标准审定、发布、版本和引用追溯纳入系统流程",
            "让标准自动驱动语义层、质检规则和数据模型",
            "形成跨部门可共享、可订阅、可审核运营的标准资产",
        ], "green", "pale_green"),
    ]
    for i, (title, items, color, fill_color) in enumerate(columns):
        x = Inches(0.8 + i * 6.15)
        shape(slide, x, Inches(1.2), Inches(5.5), Inches(5.0),
              color=fill_color, line="line")
        text(slide, x + Inches(0.32), Inches(1.55), Inches(4.8),
             Inches(0.38), title, size=22, color=color, bold=True)
        bullets(slide, x + Inches(0.32), Inches(2.25), Inches(4.85),
                Inches(2.9), items, size=17)


def design_thinking(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "设计思想：标准即单一权威源", page)
    principles = [
        ("单一权威源", "标准条款是语义、规则和模型的源头，不让下游反向污染标准"),
        ("条款级粒度", "所有引用、编辑、审定、影响分析都精确到条款和数据元"),
        ("显式派生链路", "用派生血缘记录标准到语义、规则、缺陷和模型的关系"),
        ("人机协同治理", "专家负责判断，系统负责结构化、校验、追溯和影响分析"),
        ("柔性变更", "换版后旧派生不直接删除，而是 stale，给业务平滑过渡"),
        ("渐进接入", "不要求一次替换存量系统，可先接入标准、语义或质检其中一环"),
    ]
    for i, (title, body) in enumerate(principles):
        x = Inches(0.75 + (i % 3) * 4.15)
        y = Inches(1.25 + (i // 3) * 2.18)
        shape(slide, x, y, Inches(3.55), Inches(1.42),
              color="white", line="line")
        text(slide, x + Inches(0.22), y + Inches(0.2), Inches(3.0),
             Inches(0.28), title, size=17, color="green", bold=True)
        text(slide, x + Inches(0.22), y + Inches(0.62), Inches(3.05),
             Inches(0.52), body, size=12, color="ink")


def architecture(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "总体架构：业务工作台 + 生命周期服务 + 标准资产库", page)
    layers = [
        ("客户工作台", "采集、分析、起草、审定、发布、派生、市场共享"),
        ("生命周期服务", "文档结构化、协同编辑、审定流、版本发布、影响分析"),
        ("派生引擎", "标准 → 语义层 / 质检规则 / 缺陷分类 / 数据模型"),
        ("标准资产库", "条款、术语、数据元、值域、版本、引用、派生血缘"),
        ("异步治理", "Outbox 事件、失败重试、回滚、stale 状态管理"),
        ("下游系统", "NL2SQL、数据质检、数据建模、元数据治理、跨部门共享"),
    ]
    for i, (title, desc) in enumerate(layers):
        y = Inches(1.05 + i * 0.83)
        fill_color = "pale_green" if i in (0, 2, 5) else "white"
        shape(slide, Inches(1.0), y, Inches(11.0), Inches(0.58),
              color=fill_color, line="line")
        text(slide, Inches(1.25), y + Inches(0.14), Inches(1.7),
             Inches(0.24), title, size=14, color="green", bold=True)
        text(slide, Inches(3.05), y + Inches(0.14), Inches(8.4),
             Inches(0.24), desc, size=13, color="ink")


def lifecycle_flow(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "使用流程：一份标准如何完成全生命周期管理", page)
    steps = [
        ("1", "采集", "上传标准文档、XMI 或录入企业标准"),
        ("2", "结构化", "抽取条款树、术语、数据元、值域"),
        ("3", "起草", "专家在线编辑、引用标准、锁定条款"),
        ("4", "审定", "发起审定、评论、引用核验、形成结论"),
        ("5", "发布", "生成版本快照，进入 released 状态"),
        ("6", "派生", "生成语义提示、质检规则、缺陷码、数据模型"),
        ("7", "治理", "影响分析、回滚、stale 标记、跨版本对比"),
        ("8", "共享", "市场目录、订阅、审核、组织级可见"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.65 + (i % 4) * 3.15)
        y = Inches(1.25 + (i // 4) * 2.35)
        shape(slide, x, y, Inches(2.65), Inches(1.55),
              color="white", line="line")
        text(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.48),
             Inches(0.34), num, size=20, color="green", bold=True)
        text(slide, x + Inches(0.74), y + Inches(0.2), Inches(1.85),
             Inches(0.28), title, size=17, color="ink", bold=True)
        text(slide, x + Inches(0.22), y + Inches(0.74), Inches(2.18),
             Inches(0.45), desc, size=11, color="muted")


def data_model(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "核心对象：把标准拆成可治理的数据资产", page)
    rows = [
        ("标准文档", "std_document", "标准资产入口，记录来源、编码、标题、标签和负责人"),
        ("标准版本", "std_document_version", "每次起草、审定、发布都形成可追溯版本"),
        ("条款", "std_clause", "树形条款路径，是引用、审定和影响分析的最小单位"),
        ("术语/数据元/值域", "std_term / std_data_element / std_value_domain", "承载业务语义、字段定义、取值范围和约束"),
        ("审定记录", "std_review_*", "记录审定 round、评论、引用核验和结论"),
        ("派生血缘", "std_derived_link", "追踪标准到语义层、规则、缺陷和模型的派生关系"),
        ("市场资产", "std_market_*", "支持目录、订阅、审核和组织级可见性"),
    ]
    x0, y0 = Inches(0.7), Inches(1.08)
    widths = [Inches(1.75), Inches(3.25), Inches(7.55)]
    for i, h in enumerate(["对象", "表/模块", "作用"]):
        x = x0 + sum(widths[:i])
        shape(slide, x, y0, widths[i], Inches(0.42), color="green",
              line="green", rounded=False)
        text(slide, x + Inches(0.08), y0 + Inches(0.09),
             widths[i] - Inches(0.15), Inches(0.22), h, size=11,
             color="white", bold=True)
    for r, row in enumerate(rows):
        y = y0 + Inches(0.42 + r * 0.62)
        for i, value in enumerate(row):
            x = x0 + sum(widths[:i])
            shape(slide, x, y, widths[i], Inches(0.62), color="white",
                  line="line", rounded=False)
            text(slide, x + Inches(0.08), y + Inches(0.14),
                 widths[i] - Inches(0.15), Inches(0.28), value,
                 size=11.5 if i != 2 else 11, color="ink",
                 bold=(i == 0))


def ingestion_analysis(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "阶段一：采集与结构化分析", page)
    cards = [
        ("输入来源", [
            "国家/行业/地方/企业标准文档",
            "已有 XMI 或内部标准模型",
            "互联网公开标准和人工粘贴内容",
        ], "blue"),
        ("结构化结果", [
            "条款树与章节层级",
            "术语和定义",
            "数据元、值域、约束、单位、分类",
        ], "green"),
        ("智能辅助", [
            "文档类型识别",
            "相似条款检索",
            "重复/冲突初步提示",
        ], "amber"),
    ]
    for i, (title, items, color) in enumerate(cards):
        x = Inches(0.8 + i * 4.05)
        shape(slide, x, Inches(1.35), Inches(3.45), Inches(4.55),
              color="white", line=color)
        text(slide, x + Inches(0.25), Inches(1.68), Inches(2.8),
             Inches(0.35), title, size=20, color=color, bold=True)
        bullets(slide, x + Inches(0.25), Inches(2.38), Inches(2.9),
                Inches(2.6), items, size=16)


def drafting_review(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "阶段二：专家起草与审定协同", page)
    shape(slide, Inches(0.8), Inches(1.25), Inches(5.5), Inches(4.8),
          color="white", line="green")
    text(slide, Inches(1.1), Inches(1.6), Inches(4.6), Inches(0.35),
         "起草协同", size=22, color="green", bold=True)
    bullets(slide, Inches(1.1), Inches(2.25), Inches(4.7), Inches(2.6), [
        "条款级在线编辑，避免多人覆盖同一段内容",
        "引用助手帮助专家插入可追溯引用",
        "一致性校验发现术语、数据元、值域冲突",
        "编辑过程沉淀为结构化标准资产",
    ], size=17)
    shape(slide, Inches(7.0), Inches(1.25), Inches(5.2), Inches(4.8),
          color="white", line="blue")
    text(slide, Inches(7.3), Inches(1.6), Inches(4.4), Inches(0.35),
         "审定协同", size=22, color="blue", bold=True)
    bullets(slide, Inches(7.3), Inches(2.25), Inches(4.4), Inches(2.6), [
        "发起审定 round，明确 reviewer 和审定目标",
        "对引用、条款和评论形成闭环处理",
        "发布前检查阻塞问题，避免带病发布",
        "审定记录保留为后续追溯依据",
    ], size=17)


def publish_derivation(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "阶段三：发布后自动驱动下游治理", page)
    text(slide, Inches(0.85), Inches(1.18), Inches(11.7), Inches(0.42),
         "发布不是文档归档，而是触发标准资产对下游系统的统一驱动。",
         size=20, color="green", bold=True)
    targets = [
        ("智能问数语义层", "标准术语、字段含义和值域别名进入 NL2SQL grounding"),
        ("数据质检规则", "强制项、枚举、范围、格式约束派生为 QC 规则"),
        ("缺陷分类体系", "标准约束映射到可统计、可闭环的缺陷编码"),
        ("数据建模", "生成 CDM / LDM / PDM，并输出 DDL / XMI"),
        ("影响分析", "标准升版时识别受影响的条款、数据元和下游派生"),
        ("回滚与过期", "旧派生进入 stale 状态，保留业务过渡窗口"),
    ]
    for i, (title, body) in enumerate(targets):
        x = Inches(0.82 + (i % 3) * 4.1)
        y = Inches(2.1 + (i // 3) * 1.85)
        shape(slide, x, y, Inches(3.45), Inches(1.2), color="white",
              line="line")
        text(slide, x + Inches(0.22), y + Inches(0.18), Inches(3.0),
             Inches(0.26), title, size=15.5, color="green", bold=True)
        text(slide, x + Inches(0.22), y + Inches(0.58), Inches(2.95),
             Inches(0.4), body, size=10.8, color="ink")


def governance(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "变更治理：标准升版时不再靠人工排查", page)
    flow = [
        ("标准换版", "新版本发布"),
        ("影响识别", "跨标准/跨版本影响图谱"),
        ("派生处理", "旧派生 stale，新派生 active"),
        ("业务过渡", "保留旧规则，按需回滚或切换"),
    ]
    for i, (title, desc) in enumerate(flow):
        x = Inches(0.9 + i * 3.05)
        shape(slide, x, Inches(1.55), Inches(2.45), Inches(1.18),
              color="white", line="green")
        text(slide, x + Inches(0.18), Inches(1.78), Inches(2.0),
             Inches(0.26), title, size=17, color="green", bold=True,
             align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.18), Inches(2.18), Inches(2.0),
             Inches(0.24), desc, size=11, color="muted",
             align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            conn = slide.shapes.add_connector(
                1, x + Inches(2.45), Inches(2.14),
                x + Inches(2.9), Inches(2.14)
            )
            conn.line.color.rgb = c("line")
    bullets(slide, Inches(1.0), Inches(4.0), Inches(10.8), Inches(1.4), [
        "影响图谱把标准条款、引用关系、相似条款和下游派生连接起来",
        "批量回滚允许治理人员在异常派生或错误发布后快速恢复",
        "Outbox dead-letter 运维让异步任务失败可见、可查、可重试",
    ], size=18)


def market_share(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "共享复用：让标准成为可运营资产", page)
    items = [
        ("市场目录", "已发布标准进入目录，可搜索、查看资产计数和版本信息"),
        ("版本对比", "不同版本按条款、数据元、术语、值域做结构化 diff"),
        ("订阅跟踪", "用户订阅标准后，系统提示是否存在新版本"),
        ("上架审核", "标准进入市场前经过提交、审核、通过/拒绝流程"),
        ("组织可见", "按公开、组织、私有控制跨部门/跨租户可见范围"),
    ]
    for i, (title, desc) in enumerate(items):
        y = Inches(1.12 + i * 0.92)
        shape(slide, Inches(1.0), y, Inches(11.0), Inches(0.64),
              color="white", line="line")
        text(slide, Inches(1.25), y + Inches(0.15), Inches(1.65),
             Inches(0.24), title, size=15, color="green", bold=True)
        text(slide, Inches(3.0), y + Inches(0.15), Inches(8.5),
             Inches(0.24), desc, size=13.2, color="ink")


def roles(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "角色视角：不同用户如何使用", page)
    roles_data = [
        ("标准管理员", "配置目录、审核上架、处理异常、组织共享"),
        ("标准编辑", "采集标准、起草条款、维护数据元和值域"),
        ("审定专家", "核验引用、处理评论、给出审定结论"),
        ("数据工程师", "消费 DDL / XMI / 质检规则 / 数据模型"),
        ("业务用户", "通过智能问数、目录订阅和版本提醒间接受益"),
    ]
    for i, (role, desc) in enumerate(roles_data):
        x = Inches(0.75 + (i % 2) * 6.1)
        y = Inches(1.15 + (i // 2) * 1.65)
        w = Inches(5.45)
        shape(slide, x, y, w, Inches(1.1), color="white", line="line")
        text(slide, x + Inches(0.25), y + Inches(0.18), Inches(1.7),
             Inches(0.28), role, size=16, color="green", bold=True)
        text(slide, x + Inches(1.95), y + Inches(0.18), Inches(3.0),
             Inches(0.5), desc, size=12.5, color="ink")


def scenario(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "典型场景：自然资源数据标准升版", page)
    text(slide, Inches(0.85), Inches(1.15), Inches(11.5), Inches(0.42),
         "以“地类编码/图斑面积/用途分类”等标准变化为例，系统把变更从标准层传导到数据治理层。",
         size=19, color="green", bold=True)
    bullets(slide, Inches(1.0), Inches(2.1), Inches(5.5), Inches(2.8), [
        "专家上传新标准版本并结构化为条款和数据元",
        "系统识别与旧版本之间的数据元和值域差异",
        "审定专家核验关键引用和变更说明",
        "发布后自动生成新的语义提示、质检规则和模型约束",
    ], size=17)
    bullets(slide, Inches(7.0), Inches(2.1), Inches(5.2), Inches(2.8), [
        "数据工程师查看影响图谱，确认受影响字段和规则",
        "旧规则进入 stale，业务系统可平滑切换",
        "跨部门用户订阅标准，收到新版本可见状态",
        "组织级权限保证企业/部门标准共享可控",
    ], size=17)


def value(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "客户价值：从标准管理到主动治理", page)
    values = [
        ("降低专家重复劳动", "字段、值域、规则不再反复人工解释"),
        ("提升数据质量", "标准约束直接驱动质检和缺陷闭环"),
        ("支撑智能问数", "标准术语进入 NL2SQL 语义对齐链路"),
        ("降低升版风险", "影响分析和 stale 机制让变更可控"),
        ("促进跨部门复用", "市场目录、订阅和组织可见性提升共享效率"),
        ("形成治理资产沉淀", "标准、血缘、审定记录和派生结果长期可追溯"),
    ]
    for i, (title, desc) in enumerate(values):
        x = Inches(0.82 + (i % 3) * 4.05)
        y = Inches(1.25 + (i // 3) * 2.05)
        shape(slide, x, y, Inches(3.42), Inches(1.35),
              color="white", line="line")
        text(slide, x + Inches(0.22), y + Inches(0.2), Inches(3.0),
             Inches(0.25), title, size=15.5, color="green", bold=True)
        text(slide, x + Inches(0.22), y + Inches(0.62), Inches(2.9),
             Inches(0.48), desc, size=11.5, color="ink")


def rollout(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "落地建议：先做一个标准域闭环试点", page)
    phases = [
        ("第一步", "选定标准域", "选择一个客户高频使用、变更影响明显的标准体系"),
        ("第二步", "建立标准资产", "导入标准文档，结构化条款、数据元、值域"),
        ("第三步", "跑通审定发布", "让专家在系统内完成起草、审定和发布"),
        ("第四步", "接入一个下游", "优先接 NL2SQL、质检规则或数据模型其中之一"),
        ("第五步", "形成运营机制", "建立订阅、审核、组织共享和版本变更复盘"),
    ]
    for i, (step, title, desc) in enumerate(phases):
        y = Inches(1.12 + i * 0.92)
        shape(slide, Inches(0.9), y, Inches(11.5), Inches(0.64),
              color="white", line="line")
        text(slide, Inches(1.15), y + Inches(0.15), Inches(1.1),
             Inches(0.24), step, size=13, color="muted", bold=True)
        text(slide, Inches(2.25), y + Inches(0.15), Inches(2.1),
             Inches(0.24), title, size=15, color="green", bold=True)
        text(slide, Inches(4.35), y + Inches(0.15), Inches(7.6),
             Inches(0.24), desc, size=13, color="ink")


def closing(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    header(slide, "总结：把标准变成可执行的治理底座", page)
    text(slide, Inches(1.0), Inches(1.35), Inches(10.7), Inches(0.7),
         "数据标准全生命周期智能化管理的核心，不是多一个文档库，而是建立一套“标准驱动治理”的机制。",
         size=25, color="ink", bold=True)
    bullets(slide, Inches(1.1), Inches(2.55), Inches(10.2), Inches(2.0), [
        "标准从静态文本转为结构化、版本化、可审定、可追溯的资产",
        "发布后的标准自动驱动语义、规则、缺陷分类和数据模型",
        "标准变更不再靠人工排查，而是通过影响图谱、stale 和回滚机制治理",
        "跨部门共享通过目录、订阅、审核和组织可见性形成运营闭环",
    ], size=19)
    text(slide, Inches(1.1), Inches(5.55), Inches(10.2), Inches(0.4),
         "最终目标：让标准成为数据治理和智能体能力的统一语义底座。",
         size=21, color="green", bold=True)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    title_slide(prs)
    pain_goal(prs, 2)
    design_thinking(prs, 3)
    architecture(prs, 4)
    lifecycle_flow(prs, 5)
    data_model(prs, 6)
    ingestion_analysis(prs, 7)
    drafting_review(prs, 8)
    publish_derivation(prs, 9)
    governance(prs, 10)
    market_share(prs, 11)
    roles(prs, 12)
    scenario(prs, 13)
    value(prs, 14)
    rollout(prs, 15)
    closing(prs, 16)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
