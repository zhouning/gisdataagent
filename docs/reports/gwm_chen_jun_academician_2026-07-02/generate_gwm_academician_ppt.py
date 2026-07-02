from __future__ import annotations

import shutil
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "Geospatial_World_Model_面向陈军院士交流.pptx"
PDF_PATH = OUT_DIR / "Geospatial_World_Model_面向陈军院士交流.pdf"
OUTLINE_PATH = OUT_DIR / "outline.md"
QA_PATH = OUT_DIR / "qa_report.md"

W, H = 13.333333, 7.5
FONT_CN = "PingFang SC"

plt.rcParams["font.sans-serif"] = ["Hiragino Sans GB", "Arial Unicode MS", "PingFang SC", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


COLORS = {
    "ink": rgb("18212B"),
    "muted": rgb("5F6E7C"),
    "faint": rgb("EEF3F7"),
    "line": rgb("D7E1EA"),
    "white": rgb("FFFFFF"),
    "dark": rgb("102033"),
    "blue": rgb("145D8F"),
    "teal": rgb("1E827D"),
    "green": rgb("3F8B5F"),
    "amber": rgb("A76325"),
    "red": rgb("B8453F"),
    "purple": rgb("6558A8"),
    "paper": rgb("F7F9FB"),
    "panel": rgb("FFFFFF"),
    "panel_blue": rgb("EAF3F9"),
    "panel_green": rgb("EBF6EF"),
    "panel_amber": rgb("FAF0E3"),
    "panel_red": rgb("F8EDED"),
}


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill(slide, color=COLORS["paper"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def set_font(run, size: int, bold: bool, color: RGBColor):
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 13,
    bold: bool = False,
    color: RGBColor = COLORS["ink"],
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    for idx, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        if line_spacing is not None:
            para.line_spacing = line_spacing
        for run in para.runs:
            set_font(run, size, bold, color)
    return box


def add_title(slide, title: str, kicker: str, page: int):
    add_text(slide, 0.68, 0.28, 8.4, 0.22, kicker, 8, True, COLORS["teal"])
    add_text(slide, 0.68, 0.62, 11.4, 0.48, title, 23, True, COLORS["ink"])
    add_text(slide, 12.15, 0.34, 0.5, 0.18, f"{page:02d}", 8, True, COLORS["muted"], PP_ALIGN.RIGHT)


def add_footer(slide, text: str = "Geospatial World Model / TWM 技术交流材料"):
    add_text(slide, 0.68, 7.16, 8.4, 0.16, text, 6, False, COLORS["muted"])


def add_notes(slide, notes: str):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = notes


def add_panel(slide, x, y, w, h, fill_color=COLORS["panel"], line_color=COLORS["line"], radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.color.rgb = line_color
    shp.line.width = Pt(0.8)
    return shp


def add_bar(slide, x, y, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def add_chip(slide, x, y, text, fill_color, w=None, text_color=COLORS["white"]):
    w = w or max(0.5, 0.14 * len(text) + 0.28)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    add_text(slide, x + 0.04, y + 0.055, w - 0.08, 0.14, text, 7, True, text_color, PP_ALIGN.CENTER)
    return shp


def add_card(slide, x, y, w, h, title, body, accent=COLORS["blue"], fill_color=COLORS["panel"], ts=13, bs=10):
    add_panel(slide, x, y, w, h, fill_color)
    add_bar(slide, x, y, h, accent)
    add_text(slide, x + 0.22, y + 0.16, w - 0.34, 0.25, title, ts, True, COLORS["ink"])
    add_text(slide, x + 0.22, y + 0.56, w - 0.34, h - 0.66, body, bs, False, COLORS["muted"], line_spacing=0.92)


def add_bullets(slide, x, y, w, h, bullets: Iterable[str], size=10, color=COLORS["muted"]):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        for run in para.runs:
            set_font(run, size, False, color)
    return box


def add_connector(slide, x1, y1, x2, y2, color=COLORS["line"], width=1.0, arrow=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_image_fit(slide, path: Path, x, y, w, h):
    img = Image.open(path)
    iw, ih = img.size
    ratio = min(w / iw, h / ih)
    dw, dh = iw * ratio, ih * ratio
    return slide.shapes.add_picture(str(path), Inches(x + (w - dw) / 2), Inches(y + (h - dh) / 2), Inches(dw), Inches(dh))


def make_cover(src: Path, dst: Path):
    img = Image.open(src).convert("RGB")
    img = ImageOps.autocontrast(img)
    overlay = Image.new("RGB", img.size, (7, 18, 30))
    img = Image.blend(img, overlay, 0.54)
    img.save(dst)


def ensure_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    selected = {
        "cover_source": ROOT / "tests/e2e/screenshots/twm_demo_workflow.png",
        "twm_arch": ROOT / "docs/assets/twm_architecture_overview.png",
        "twm_upstream": ROOT / "docs/assets/twm_upstream_foundation_overview.png",
        "twm_plan": ROOT / "tests/e2e/screenshots/twm_overview_plan.png",
        "twm_evidence": ROOT / "tests/e2e/screenshots/twm_data_evidence.png",
        "flus_positioning": ROOT / "docs/assets/twm_flus_positioning_boundary.png",
    }
    for key, src in selected.items():
        if src.exists():
            shutil.copyfile(src, ASSET_DIR / src.name)
    make_cover(selected["cover_source"], ASSET_DIR / "cover_gwm_academician.png")


def make_evidence_chart() -> Path:
    groups = ["Direct CA\nfixed baseline", "ANN+CA\npaired 85 cases"]
    twm = np.array([0.195662, 0.073243])
    flus = np.array([0.150955, 0.144175])
    x = np.arange(len(groups))
    width = 0.32
    fig, ax = plt.subplots(figsize=(7.2, 3.8), dpi=180)
    ax.bar(x - width / 2, twm, width, label="TWM", color="#1E827D")
    ax.bar(x + width / 2, flus, width, label="FLUS", color="#A76325")
    ax.set_ylim(0, 0.24)
    ax.set_ylabel("Change FoM")
    ax.set_xticks(x)
    ax.set_xticklabels(groups)
    ax.grid(axis="y", alpha=0.22)
    ax.legend(frameon=False, loc="upper right")
    for i, (a, b) in enumerate(zip(twm, flus, strict=False)):
        ax.text(i - width / 2, a + 0.006, f"{a:.3f}", ha="center", va="bottom", fontsize=8)
        ax.text(i + width / 2, b + 0.006, f"{b:.3f}", ha="center", va="bottom", fontsize=8)
    ax.spines[["top", "right", "left"]].set_visible(False)
    fig.tight_layout()
    path = ASSET_DIR / "twm_flus_evidence_boundary_chart.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def make_route_chart() -> Path:
    labels = ["GIS 表达", "数字孪生", "GeoAI", "土地模拟", "GWM"]
    state = np.array([0.95, 0.90, 0.65, 0.70, 0.88])
    action = np.array([0.25, 0.35, 0.32, 0.52, 0.86])
    evidence = np.array([0.65, 0.62, 0.35, 0.38, 0.82])
    x = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(8.0, 4.0), dpi=180)
    ax.plot(x, state, marker="o", lw=2.4, color="#145D8F", label="state representation")
    ax.plot(x, action, marker="o", lw=2.4, color="#1E827D", label="action consequence")
    ax.plot(x, evidence, marker="o", lw=2.4, color="#A76325", label="evidence boundary")
    ax.set_ylim(0, 1.0)
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.grid(axis="y", alpha=0.22)
    ax.legend(frameon=False, loc="lower right")
    ax.spines[["top", "right", "left"]].set_visible(False)
    fig.tight_layout()
    path = ASSET_DIR / "gwm_route_capability_chart.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def section_slide(prs: Presentation, page: int, section: str, title: str, subtitle: str, color=COLORS["dark"]):
    slide = blank(prs)
    fill(slide, color)
    add_chip(slide, 0.72, 0.62, section, COLORS["teal"], w=1.18)
    add_text(slide, 0.72, 1.42, 10.8, 0.7, title, 30, True, COLORS["white"])
    add_text(slide, 0.76, 2.38, 9.6, 0.7, subtitle, 15, False, rgb("D8E2EA"), line_spacing=1.0)
    add_text(slide, 12.04, 6.95, 0.65, 0.2, f"{page:02d}", 8, True, rgb("D8E2EA"), PP_ALIGN.RIGHT)
    return slide


def build_deck() -> Path:
    ensure_assets()
    evidence_chart = make_evidence_chart()
    route_chart = make_route_chart()

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    outline: list[str] = []

    # 1
    slide = blank(prs)
    slide.shapes.add_picture(str(ASSET_DIR / "cover_gwm_academician.png"), 0, 0, width=Inches(W), height=Inches(H))
    add_chip(slide, 0.72, 0.64, "院士交流材料", COLORS["teal"], w=1.45)
    add_text(slide, 0.72, 1.36, 9.4, 0.86, "Geospatial World Model：\n面向自然资源治理的地理空间世界模型探索", 26, True, COLORS["white"])
    add_text(slide, 0.76, 3.12, 7.8, 0.48, "以 TWM 作为自然资源行业实例，讨论 GIScience、GeoAI 与行动后果推演的交叉方向。", 14, False, rgb("E6EEF5"))
    add_text(slide, 0.76, 6.82, 4.6, 0.22, "GIS Data Agent / 2026-07-02", 8, False, rgb("D8E2EA"))
    add_notes(slide, "开场强调：这不是产品汇报，而是希望就 Geospatial World Model 的科学定位、边界和研究路线请专家判断。")
    outline.append("标题页")

    # 2
    slide = blank(prs)
    fill(slide)
    add_title(slide, "这份材料的交流目标", "面向院士交流的定位", 2)
    add_text(slide, 0.92, 1.28, 11.4, 0.4, "不是做一个系统演示，而是请判断一个研究方向是否成立。", 21, True, COLORS["blue"], PP_ALIGN.CENTER)
    items = [
        ("概念是否成立", "Geospatial World Model 是否能作为 GIScience + GeoAI 的下一类研究对象？"),
        ("边界是否严谨", "它和 GIS、数字孪生、土地利用模拟、空间优化、GeoAI foundation model 的边界如何划分？"),
        ("路线是否值得推进", "自然资源 TWM 是否适合作为 GWM 的第一类行业实例与验证场景？"),
    ]
    for i, (t, b) in enumerate(items):
        add_card(slide, 0.95 + i * 4.1, 2.28, 3.45, 2.05, t, b, [COLORS["blue"], COLORS["teal"], COLORS["amber"]][i], [COLORS["panel_blue"], COLORS["panel_green"], COLORS["panel_amber"]][i], 14, 11)
    add_text(slide, 1.15, 5.52, 11.0, 0.34, "因此，后续会主动说明已有工作、实验边界和当前不足，避免把概念包装成已完成结论。", 15, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("交流目标")

    # 3
    slide = blank(prs)
    fill(slide)
    add_title(slide, "从 GIS / GeoAI 到 GWM：问题正在变化", "研究背景", 3)
    add_image_fit(slide, route_chart, 0.78, 1.25, 6.4, 3.4)
    add_card(slide, 7.75, 1.22, 4.65, 1.08, "GIS 的强项", "表达、组织、查询、分析空间信息。", COLORS["blue"], COLORS["panel_blue"], 14, 11)
    add_card(slide, 7.75, 2.66, 4.65, 1.08, "GeoAI 的强项", "识别、表征、预测遥感与空间数据模式。", COLORS["teal"], COLORS["panel_green"], 14, 11)
    add_card(slide, 7.75, 4.10, 4.65, 1.08, "GWM 的目标", "推演行动后果，约束模型结论，服务可审计决策。", COLORS["amber"], COLORS["panel_amber"], 14, 11)
    add_text(slide, 0.9, 5.75, 11.45, 0.36, "从“描述世界”到“识别世界”之后，自然资源治理更需要“推演世界”。", 17, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("GIS/GeoAI 到 GWM")

    # 4
    slide = blank(prs)
    fill(slide)
    add_title(slide, "World Model 的最低共识", "概念基线", 4)
    add_text(slide, 0.88, 1.25, 11.5, 0.55, "世界模型不是“生成好看的未来画面”，而是学习环境状态、演化规律和行动后果。", 22, True, COLORS["blue"], PP_ALIGN.CENTER)
    steps = [("State", "世界现在是什么状态"), ("Dynamics", "世界如何随时间变化"), ("Action", "行动如何改变未来"), ("Planning", "如何基于推演选择方案")]
    for i, (t, b) in enumerate(steps):
        x = 0.95 + i * 3.0
        add_panel(slide, x, 2.55, 2.25, 1.5, COLORS["panel"])
        add_chip(slide, x + 0.72, 2.85, t, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"]][i], w=0.78)
        add_text(slide, x + 0.18, 3.42, 1.9, 0.25, b, 10, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < 3:
            add_connector(slide, x + 2.28, 3.28, x + 2.78, 3.28, COLORS["line"], 1.2, True)
    add_panel(slide, 1.25, 5.05, 10.85, 0.78, COLORS["panel_blue"])
    add_text(slide, 1.5, 5.31, 10.35, 0.2, "普通形式：p(S_{t+1} | S_t, A_t)；GWM 需要把 S 和 A 都换成地理空间与治理语义。", 15, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide, "References: Ha & Schmidhuber 2018; Dreamer 系列；model-based planning literature.")
    outline.append("World Model 共识")

    # 5
    slide = blank(prs)
    fill(slide)
    add_title(slide, "为什么 GWM 不是普通 World Model 的搬运", "地理空间特殊性", 5)
    axes = [
        ("对象-场二元", "地块/道路/行政区 + 遥感/栅格场"),
        ("拓扑关系", "相邻、包含、连通、缓冲、上下游"),
        ("尺度效应", "parcel / township / county 改变结论"),
        ("空间异质", "规则、地貌、经济与数据质量区域差异"),
        ("治理规则", "红线、基本农田、用途管制、审批要求"),
        ("证据边界", "来源、时间、CRS、置信度、法律效力"),
    ]
    for i, (t, b) in enumerate(axes):
        x = 0.82 + (i % 3) * 4.1
        y = 1.35 + (i // 3) * 2.05
        add_card(slide, x, y, 3.45, 1.45, t, b, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["red"], COLORS["purple"]][i], COLORS["panel"], 13, 10)
    add_text(slide, 0.95, 5.86, 11.4, 0.34, "GWM 的核心难点不是“再训练一个模型”，而是把空间、规则、证据和行动统一进模型状态。", 16, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("GWM 特殊性")

    # 6
    slide = blank(prs)
    fill(slide)
    add_title(slide, "Geospatial World Model 的形式化定义", "核心定义", 6)
    add_panel(slide, 0.9, 1.28, 11.55, 1.05, COLORS["dark"])
    add_text(slide, 1.15, 1.66, 11.05, 0.18, "GWM = <S_G, A_G, R_G, E_G, T_G, U_G, C_G, M_G>", 20, True, COLORS["white"], PP_ALIGN.CENTER)
    defs = [
        ("S_G", "geospatial state"),
        ("A_G", "geospatial action"),
        ("R_G", "rules / constraints"),
        ("E_G", "evidence / uncertainty"),
        ("T_G", "transition dynamics"),
        ("U_G", "utility / regret"),
        ("C_G", "claim boundary"),
        ("M_G", "failure memory"),
    ]
    for i, (sym, desc) in enumerate(defs):
        x = 0.85 + (i % 4) * 3.1
        y = 2.9 + (i // 4) * 1.25
        add_panel(slide, x, y, 2.55, 0.82, COLORS["panel"])
        add_text(slide, x + 0.12, y + 0.16, 0.58, 0.2, sym, 13, True, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["red"], COLORS["purple"], COLORS["blue"], COLORS["teal"]][i], PP_ALIGN.CENTER)
        add_text(slide, x + 0.75, y + 0.18, 1.6, 0.18, desc, 9, False, COLORS["muted"])
    add_text(slide, 1.0, 5.92, 11.2, 0.33, "定义重点：GWM 不是普通预测器，而是受空间、证据和治理约束的行动世界模型。", 16, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("GWM 形式化定义")

    # 7
    slide = blank(prs)
    fill(slide)
    add_title(slide, "与已有路线的边界：不能把相邻工作说成空白", "学术定位", 7)
    cols = [0.68, 2.05, 4.4, 7.0, 9.65]
    widths = [1.15, 2.05, 2.25, 2.35, 2.85]
    headers = ["路线", "主要对象", "核心输出", "强项", "GWM 的差异"]
    for x, w, htxt in zip(cols, widths, headers, strict=False):
        add_panel(slide, x, 1.2, w, 0.42, COLORS["dark"], radius=False)
        add_text(slide, x + 0.04, 1.32, w - 0.08, 0.12, htxt, 7, True, COLORS["white"], PP_ALIGN.CENTER)
    rows = [
        ("GIS", "空间数据与关系", "图层/分析/制图", "表达严谨", "加入未来与行动后果"),
        ("数字孪生", "实时状态", "同步/监控/可视", "状态可见", "从同步走向可审计推演"),
        ("GeoAI/FM", "遥感与空间模式", "识别/表征/预测", "模型能力强", "加入规则、行动、证据边界"),
        ("LULC 模拟", "土地类型转移", "未来土地利用图", "成熟基线", "从地类转移扩展到治理行动"),
        ("空间优化", "目标与约束", "方案排序", "决策导向", "优化消费 simulator，而非替代模型"),
    ]
    for i, row in enumerate(rows):
        y = 1.72 + i * 0.78
        fill_color = COLORS["panel"] if i % 2 == 0 else rgb("F3F7FA")
        for x, w, cell in zip(cols, widths, row, strict=False):
            add_panel(slide, x, y, w, 0.58, fill_color, radius=False)
            add_text(slide, x + 0.05, y + 0.16, w - 0.1, 0.13, cell, 7, False, COLORS["ink"] if x == cols[0] else COLORS["muted"])
    add_text(slide, 0.9, 6.18, 11.55, 0.25, "稳妥表述：GWM 不是“首次做地理空间模拟”，而是把状态、行动、规则、证据和规划闭环系统化。", 13, True, COLORS["red"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("相关路线边界")

    # 8
    slide = section_slide(prs, 8, "PART 1", "自然资源治理是 GWM 的高价值起点", "因为它天然同时包含空间对象、规则约束、权威证据、历史演化和行动后果。", COLORS["blue"])
    outline.append("章节：自然资源场景")

    # 9
    slide = blank(prs)
    fill(slide)
    add_title(slide, "自然资源问题为什么适合第一个实例", "行业问题结构", 9)
    rows = [
        ("空间对象明确", "地块、图斑、行政区、规划区、项目范围、保护边界"),
        ("规则约束强", "生态红线、基本农田、城镇开发边界、用途管制"),
        ("证据体系清晰", "遥感变化、权威底图、审批记录、巡查证据、政策文本"),
        ("行动后果重要", "批准、调整、保护、整治、开发都会改变未来状态和风险"),
    ]
    for i, (t, b) in enumerate(rows):
        y = 1.22 + i * 1.1
        add_chip(slide, 0.88, y + 0.12, f"{i + 1}", [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"]][i], w=0.42)
        add_panel(slide, 1.55, y, 10.85, 0.78, COLORS["panel"])
        add_bar(slide, 1.55, y, 0.78, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"]][i])
        add_text(slide, 1.82, y + 0.15, 2.4, 0.2, t, 12, True, COLORS["ink"])
        add_text(slide, 4.3, y + 0.16, 7.65, 0.18, b, 10, False, COLORS["muted"])
    add_panel(slide, 1.15, 6.0, 11.0, 0.42, COLORS["panel_green"])
    add_text(slide, 1.35, 6.13, 10.6, 0.12, "自然资源 GWM 的核心问题：在规则与证据约束下，某个空间行动会导致怎样的未来状态、风险和效用？", 13, True, COLORS["green"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("自然资源适配性")

    # 10
    slide = blank(prs)
    fill(slide)
    add_title(slide, "TWM：自然资源领域的 Geospatial World Model 实例", "系统实例", 10)
    if (ASSET_DIR / "twm_architecture_overview.png").exists():
        add_image_fit(slide, ASSET_DIR / "twm_architecture_overview.png", 0.78, 1.1, 5.8, 4.4)
    add_card(slide, 7.05, 1.2, 5.25, 1.0, "Renderer", "GIS-operational rendering：地图、证据、风险、审计视图。", COLORS["blue"], COLORS["panel_blue"], 14, 10)
    add_card(slide, 7.05, 2.55, 5.25, 1.0, "Simulator", "学术核心：行动条件下的状态、约束、效用和不确定性推演。", COLORS["green"], COLORS["panel_green"], 14, 10)
    add_card(slide, 7.05, 3.9, 5.25, 1.0, "Planner", "消费 simulator 输出，做候选方案排序、风险复核和任务生成。", COLORS["amber"], COLORS["panel_amber"], 14, 10)
    add_text(slide, 0.9, 5.86, 11.45, 0.35, "给院士交流时应强调：TWM 不是 GWM 全部，而是自然资源治理场景下的一个验证实例。", 16, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("TWM 实例")

    # 11
    slide = blank(prs)
    fill(slide)
    add_title(slide, "Governance-bound Geospatial State", "TWM 的状态表达", 11)
    add_panel(slide, 1.0, 1.25, 11.25, 0.82, COLORS["dark"])
    add_text(slide, 1.25, 1.55, 10.75, 0.16, "Geospatial State = Geometry + Topology + Semantics + Governance + Evidence + Time", 17, True, COLORS["white"], PP_ALIGN.CENTER)
    parts = [
        ("Geometry", "parcel / polygon / boundary"),
        ("Topology", "adjacent / contains / overlaps"),
        ("Semantics", "land type / use / entity"),
        ("Governance", "redline / PBF / rule hit"),
        ("Evidence", "source / time / CRS / confidence"),
        ("Time", "history / horizon / rule version"),
    ]
    for i, (t, b) in enumerate(parts):
        x = 0.85 + i * 2.05
        add_panel(slide, x, 2.8, 1.55, 1.25, COLORS["panel"])
        add_text(slide, x + 0.08, 3.05, 1.38, 0.18, t, 10, True, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["purple"], COLORS["red"]][i], PP_ALIGN.CENTER)
        add_text(slide, x + 0.1, 3.48, 1.35, 0.26, b, 7, False, COLORS["muted"], PP_ALIGN.CENTER)
    add_panel(slide, 1.2, 5.0, 10.9, 0.72, COLORS["panel_blue"])
    add_text(slide, 1.42, 5.24, 10.45, 0.18, "状态空间本身就是创新焦点：它不是 flat feature vector，而是可审计、可追溯、可跨尺度的 GIS 治理状态。", 14, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("状态表达")

    # 12
    slide = blank(prs)
    fill(slide)
    add_title(slide, "TWM Simulator：从变化预测到行动后果推演", "算法核心", 12)
    chain = [
        ("历史状态对", "S_t -> S_t+1"),
        ("Transition replay", "训练期回放"),
        ("Topology support", "前沿/邻域支撑"),
        ("Action mask", "规则可行性"),
        ("Demand / scenario", "情景约束"),
        ("Risk & utility", "风险/收益/不确定性"),
    ]
    for i, (t, b) in enumerate(chain):
        x = 0.52 + i * 2.12
        add_panel(slide, x, 2.0, 1.68, 1.25, COLORS["panel"])
        add_text(slide, x + 0.1, 2.2, 1.48, 0.18, t, 8, True, COLORS["ink"], PP_ALIGN.CENTER)
        add_text(slide, x + 0.1, 2.67, 1.48, 0.18, b, 7, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < len(chain) - 1:
            add_connector(slide, x + 1.72, 2.62, x + 2.02, 2.62, COLORS["line"], 1.2, True)
    add_panel(slide, 1.0, 4.35, 11.3, 0.92, COLORS["panel_green"])
    add_text(slide, 1.25, 4.64, 10.8, 0.25, "p(future_state, constraint_risk, utility_delta, uncertainty | current_gis_state, action, rules, evidence)", 15, True, COLORS["green"], PP_ALIGN.CENTER)
    add_text(slide, 1.0, 5.82, 11.3, 0.3, "这部分才是与土地利用模拟、GeoAI 预测模型真正拉开问题定义的地方。", 15, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("Simulator 算法核心")

    # 13
    slide = blank(prs)
    fill(slide)
    add_title(slide, "Evidence-bounded Claim：让模型结论受证据约束", "治理可信性", 13)
    ladder = [
        ("diagnostic only", "仅诊断，不建议行动"),
        ("replay ready", "历史回放通过"),
        ("pilot review", "可控试点评审"),
        ("promotion review", "候选晋级评审"),
        ("production", "需权威验证后才可进入"),
    ]
    for i, (t, b) in enumerate(ladder):
        x = 0.9 + i * 2.45
        add_panel(slide, x, 2.0, 1.9, 1.25, [COLORS["panel"], COLORS["panel_blue"], COLORS["panel_green"], COLORS["panel_amber"], COLORS["panel_red"]][i])
        add_text(slide, x + 0.08, 2.22, 1.72, 0.18, t, 8, True, COLORS["ink"], PP_ALIGN.CENTER)
        add_text(slide, x + 0.12, 2.7, 1.66, 0.22, b, 7, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < len(ladder) - 1:
            add_connector(slide, x + 1.95, 2.6, x + 2.31, 2.6, COLORS["line"], 1.0, True)
    add_card(slide, 1.0, 4.3, 5.25, 1.1, "核心原则", "Claim(model) <= Evidence(state, replay, canary, rollback, registry)", COLORS["blue"], COLORS["panel_blue"], 14, 12)
    add_card(slide, 7.0, 4.3, 5.25, 1.1, "不是削弱模型", "把模型从“大胆输出”约束为“可审计、可回放、可回滚”。", COLORS["green"], COLORS["panel_green"], 14, 11)
    add_footer(slide)
    outline.append("证据边界")

    # 14
    slide = blank(prs)
    fill(slide)
    add_title(slide, "与 GeoSOS-FLUS 的关系：强基线，而不是被简单替代", "对比边界", 14)
    add_card(slide, 0.85, 1.28, 5.55, 3.5, "GeoSOS / FLUS", "成熟土地利用变化模拟与空间优化路线。\n\n核心对象：land-use transition and allocation。\n\n优势：ANN suitability + CA allocation 机制成熟，整体地图模拟稳定。", COLORS["amber"], COLORS["panel_amber"], 15, 11)
    add_card(slide, 6.95, 1.28, 5.55, 3.5, "TWM / GWM", "面向自然资源治理的行动世界模型。\n\n核心对象：action-conditioned territorial governance state。\n\n目标：未来状态、约束风险、规划效用、不确定性和证据边界。", COLORS["green"], COLORS["panel_green"], 15, 11)
    add_text(slide, 1.05, 5.65, 11.1, 0.35, "严谨说法：FLUS 可作为土地变化模拟强基线；TWM 需要证明更宽的治理世界模型问题，而不是只宣称替代 FLUS。", 16, True, COLORS["red"], PP_ALIGN.CENTER)
    add_footer(slide, "Source basis: GeoSOS/FLUS documentation and local TWM-vs-FLUS academic positioning notes.")
    outline.append("TWM 与 FLUS 边界")

    # 15
    slide = blank(prs)
    fill(slide)
    add_title(slide, "当前实验能支持什么，不能支持什么", "证据边界示例", 15)
    add_image_fit(slide, evidence_chart, 0.85, 1.28, 6.35, 3.55)
    add_card(slide, 7.75, 1.2, 4.65, 1.1, "能支持", "在 direct CA fixed baseline 下，TWM 的变化定位指标显示优势。", COLORS["green"], COLORS["panel_green"], 14, 11)
    add_card(slide, 7.75, 2.65, 4.65, 1.1, "不能支持", "不能据此声称 TWM 全面超过 GeoSOS-FLUS 或完整 ANN+CA 工作流。", COLORS["red"], COLORS["panel_red"], 14, 11)
    add_card(slide, 7.75, 4.1, 4.65, 1.1, "最新严谨边界", "ANN-trained FLUS 85 成功配对案例中，FLUS 在变化空间准确性上仍领先；TWM 优势主要是需求守恒。", COLORS["amber"], COLORS["panel_amber"], 14, 10)
    add_text(slide, 0.95, 5.84, 11.5, 0.34, "给院士版不应隐藏这个结果：它说明 TWM 还有真正算法突破空间，也避免概念先行。", 15, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("实验证据边界")

    # 16
    slide = blank(prs)
    fill(slide)
    add_title(slide, "真正需要突破的科学问题", "当前不足", 16)
    problems = [
        ("尺度泛化", "MAUP / FoV / 行政尺度变化下结论稳定性"),
        ("空间因果", "区分规划行动效果与空间背景相关性"),
        ("规则版本", "政策、红线、用途管制变化如何进入模型"),
        ("不确定性", "风险、证据缺口和模型置信度如何传播"),
        ("评价体系", "从 change accuracy 走向 decision regret 和 claim safety"),
        ("权威数据", "与自然资源权威底图、审批、巡查证据连接"),
    ]
    for i, (t, b) in enumerate(problems):
        x = 0.82 + (i % 3) * 4.1
        y = 1.35 + (i // 3) * 2.05
        add_card(slide, x, y, 3.45, 1.45, t, b, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["red"], COLORS["purple"]][i], COLORS["panel"], 13, 10)
    add_text(slide, 0.95, 5.86, 11.4, 0.34, "这些问题解决后，TWM 才能从工程原型走向可发表、可验证、可评审的 GWM 研究框架。", 16, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("科学问题")

    # 17
    slide = blank(prs)
    fill(slide)
    add_title(slide, "下一步：从工程原型走向学术范式", "研究路线", 17)
    roadmap = [
        ("定义", "形成 GWM 状态、行动、规则、证据、结论边界的形式化定义"),
        ("数据", "建设自然资源治理时空样本库：状态对、行动、规则、证据、结果"),
        ("模型", "实现 action-conditioned multi-head dynamics，而非单一 land-use forecast"),
        ("验证", "构建分层 benchmark：未来状态、反事实、规划收益、证据安全"),
        ("合作", "引入权威数据、专家规则和真实评审案例，形成可复核示范"),
    ]
    for i, (t, b) in enumerate(roadmap):
        y = 1.2 + i * 0.96
        add_chip(slide, 0.85, y + 0.12, f"{i + 1}", [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i], w=0.42)
        add_panel(slide, 1.5, y, 10.9, 0.68, COLORS["panel"])
        add_bar(slide, 1.5, y, 0.68, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i])
        add_text(slide, 1.78, y + 0.13, 1.2, 0.16, t, 11, True, COLORS["ink"])
        add_text(slide, 3.1, y + 0.14, 8.8, 0.14, b, 9, False, COLORS["muted"])
    add_footer(slide)
    outline.append("研究路线")

    # 18
    slide = blank(prs)
    fill(slide, COLORS["dark"])
    add_text(slide, 0.82, 0.82, 10.8, 0.52, "希望请陈军院士重点指正的问题", 27, True, COLORS["white"])
    questions = [
        ("1", "Geospatial World Model 这个概念是否能成为 GIScience / GeoAI 的严谨研究方向？"),
        ("2", "自然资源治理是否适合作为第一类 GWM 验证场景？还缺哪些关键场景？"),
        ("3", "与 GeoSOS/FLUS、数字孪生、空间优化的学术边界如何表述更稳妥？"),
        ("4", "评价体系应如何从地图精度升级为治理行动后果、证据安全和决策后悔？"),
    ]
    for i, (n, q) in enumerate(questions):
        y = 1.85 + i * 1.05
        add_chip(slide, 1.0, y, n, COLORS["teal"], w=0.42)
        add_text(slide, 1.65, y + 0.045, 10.7, 0.22, q, 14, True, rgb("E6EEF5"))
    add_footer(slide, "References: World Models; GeoAI foundation model literature; GeoSOS/FLUS; GIS Data Agent TWM local reports.")
    outline.append("请教问题")

    # 19
    slide = blank(prs)
    fill(slide)
    add_title(slide, "三个结论", "总结", 19)
    takeaways = [
        ("GWM 的定位", "让 GIS 从“空间信息表达系统”扩展为“空间行动后果推演系统”。"),
        ("TWM 的定位", "TWM 是自然资源治理场景下的 GWM 实例，不等同于 GWM 全部。"),
        ("当前边界", "概念与工程雏形已经形成；算法、权威数据、评价体系和基线验证仍需继续突破。"),
    ]
    for i, (t, b) in enumerate(takeaways):
        add_card(slide, 0.95 + i * 4.1, 1.8, 3.45, 2.35, t, b, [COLORS["blue"], COLORS["green"], COLORS["amber"]][i], [COLORS["panel_blue"], COLORS["panel_green"], COLORS["panel_amber"]][i], 15, 12)
    add_panel(slide, 1.15, 5.25, 11.0, 0.76, COLORS["dark"])
    add_text(slide, 1.42, 5.52, 10.45, 0.18, "最终目标：可预测、可反事实、可规划、可审计的地理空间治理世界模型。", 16, True, COLORS["white"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("总结")

    # 20
    slide = blank(prs)
    fill(slide)
    add_title(slide, "参考依据", "References / Source Basis", 20)
    refs = [
        "Ha & Schmidhuber, World Models, 2018.",
        "Hafner et al., Dreamer / DreamerV3 world model literature.",
        "GeoAI foundation model and Earth observation foundation model literature.",
        "GeoSOS / FLUS: land-use simulation and spatial optimization model family; Liu et al. 2017 FLUS paper.",
        "GIS Data Agent local reports: TWM theoretical foundations, novelty boundary, TWM-vs-GeoSOS/FLUS academic positioning, ANN-trained FLUS comparison.",
    ]
    add_bullets(slide, 1.0, 1.35, 11.3, 3.4, refs, 13, COLORS["ink"])
    add_panel(slide, 1.0, 5.35, 11.3, 0.76, COLORS["panel_blue"])
    add_text(slide, 1.25, 5.62, 10.8, 0.18, "本材料采用保守表述：强调研究问题和验证路线，不把工程原型包装成已完成理论。", 14, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    outline.append("参考依据")

    for s in prs.slides:
        for shape in s.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.font.name:
                            run.font.name = FONT_CN

    prs.save(PPTX_PATH)
    OUTLINE_PATH.write_text("\n".join(f"{i + 1}. {name}" for i, name in enumerate(outline)) + "\n", encoding="utf-8")
    QA_PATH.write_text(
        "\n".join(
            [
                "# QA Report",
                "",
                f"- PPTX: `{PPTX_PATH.name}`",
                f"- Slide count: {len(outline)}",
                "- Audience: Chen Jun academician / senior GIScience and geospatial information audience.",
                "- Positioning: academic and rigorous; not product marketing.",
                "",
                "## Content Checks",
                "",
                "- Defines GWM before introducing TWM.",
                "- Positions TWM as a natural-resource GWM instance, not the whole GWM category.",
                "- Separates GIS, digital twin, GeoAI, land-use simulation, spatial optimization and GWM.",
                "- Treats GeoSOS-FLUS as a strong baseline and avoids broad superiority claims.",
                "- Includes current evidence limits: direct CA advantage but ANN-trained FLUS still leads on paired change accuracy.",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    return PPTX_PATH


if __name__ == "__main__":
    print(build_deck())
