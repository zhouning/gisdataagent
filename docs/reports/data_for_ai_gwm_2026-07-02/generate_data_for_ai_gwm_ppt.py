from __future__ import annotations

import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "Data_for_AI_从时空数据治理到地理空间世界模型.pptx"
QA_PATH = OUT_DIR / "qa_report.md"
OUTLINE_PATH = OUT_DIR / "outline.md"

W = 13.333333
H = 7.5
FONT_CN = "PingFang SC"
FONT_LATIN = "Aptos"


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


COLORS = {
    "ink": rgb("17202A"),
    "muted": rgb("5F6B77"),
    "soft": rgb("F6F8FA"),
    "panel": rgb("EDF4F8"),
    "panel_green": rgb("EEF7F0"),
    "panel_amber": rgb("FAF3E7"),
    "panel_red": rgb("F8ECEB"),
    "blue": rgb("1B5C8C"),
    "teal": rgb("247C7A"),
    "green": rgb("3E8A5E"),
    "amber": rgb("B56A2A"),
    "red": rgb("B7433D"),
    "purple": rgb("675AA7"),
    "line": rgb("D7E0E8"),
    "dark": rgb("111827"),
    "white": rgb("FFFFFF"),
}


@dataclass
class DeckAsset:
    key: str
    source: Path
    local: Path
    note: str


ASSETS: list[DeckAsset] = []


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_background(slide, color=COLORS["soft"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_run_font(run, size: int, bold: bool, color: RGBColor, font: str = FONT_CN):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = COLORS["ink"],
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = FONT_CN,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for run in p.runs:
            set_run_font(run, size, bold, color, font)
    return box


def add_title(slide, title: str, kicker: str = "", page: int | None = None):
    if kicker:
        add_text(slide, 0.62, 0.26, 9.7, 0.25, kicker, 8, True, COLORS["teal"])
        y = 0.52
    else:
        y = 0.36
    add_text(slide, 0.62, y, 11.2, 0.48, title, 23, True, COLORS["ink"])
    if page is not None:
        add_text(slide, 12.2, 0.32, 0.55, 0.22, f"{page:02d}", 8, True, COLORS["muted"], PP_ALIGN.RIGHT)


def add_footer(slide, text: str = "GIS Data Agent / Data for AI 技术交流"):
    add_text(slide, 0.62, 7.17, 7.7, 0.16, text, 6, False, COLORS["muted"])


def add_notes(slide, notes: str):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = notes


def add_panel(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["line"], radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    return shp


def add_chip(slide, x, y, text, fill, color=COLORS["white"], w: float | None = None):
    w = w or max(0.75, 0.12 * len(text) + 0.35)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    add_text(slide, x + 0.06, y + 0.055, w - 0.12, 0.16, text, 8, True, color, PP_ALIGN.CENTER)
    return shp


def add_bullets(slide, x, y, w, h, bullets: Iterable[str], size=12, color=COLORS["ink"]):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            set_run_font(run, size, False, color)
    return box


def add_card(slide, x, y, w, h, title: str, body: str, fill=COLORS["white"], accent=COLORS["blue"], title_size=14, body_size=11):
    add_panel(slide, x, y, w, h, fill)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = accent
    shp.line.fill.background()
    add_text(slide, x + 0.22, y + 0.18, w - 0.35, 0.32, title, title_size, True, COLORS["ink"])
    add_text(slide, x + 0.22, y + 0.62, w - 0.35, h - 0.74, body, body_size, False, COLORS["muted"], line_spacing=0.95)


def add_connector(slide, x1, y1, x2, y2, color=COLORS["line"], width=1.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x1, y1, x2, y2, color=COLORS["line"], width=1.2):
    line = add_connector(slide, x1, y1, x2, y2, color, width)
    line.line.end_arrowhead = True
    return line


def add_image_fit(slide, path: Path, x, y, w, h):
    img = Image.open(path)
    iw, ih = img.size
    ratio = min(w / iw, h / ih)
    dw = iw * ratio
    dh = ih * ratio
    return slide.shapes.add_picture(str(path), Inches(x + (w - dw) / 2), Inches(y + (h - dh) / 2), Inches(dw), Inches(dh))


def make_cover_image(src: Path, dst: Path):
    img = Image.open(src).convert("RGB")
    img = ImageOps.autocontrast(img)
    overlay = Image.new("RGB", img.size, (10, 18, 27))
    img = Image.blend(img, overlay, 0.48)
    img.save(dst)


def ensure_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    selected = {
        "cover_source": ROOT / "tests/e2e/screenshots/twm_demo_workflow.png",
        "mmfe_quality": ROOT / "tests/e2e/screenshots/mmfe_fusion_quality_e2e.png",
        "twm_workflow": ROOT / "tests/e2e/screenshots/twm_demo_workflow.png",
        "twm_plan": ROOT / "tests/e2e/screenshots/twm_overview_plan.png",
        "twm_arch": ROOT / "docs/assets/twm_architecture_overview.png",
        "twm_upstream": ROOT / "docs/assets/twm_upstream_foundation_overview.png",
        "twm_metrics": ROOT / "docs/assets/twm_flus_v24_key_metric_comparison.png",
    }
    for key, src in selected.items():
        if src.exists():
            dst = ASSET_DIR / src.name
            shutil.copyfile(src, dst)
            ASSETS.append(DeckAsset(key, src, dst, "project asset"))
    cover = ASSET_DIR / "cover_data_for_ai.png"
    make_cover_image(selected["cover_source"], cover)
    ASSETS.append(DeckAsset("cover", selected["cover_source"], cover, "darkened cover screenshot"))


def make_metric_chart():
    labels = ["Change FoM", "Change F1", "OA", "Kappa", "Macro-F1"]
    twm = np.array([0.195662, 0.324579, 0.900997, 0.770697, 0.484675])
    flus = np.array([0.150955, 0.254339, 0.918396, 0.810473, 0.505526])
    x = np.arange(len(labels))
    width = 0.36
    fig, ax = plt.subplots(figsize=(9.5, 4.5), dpi=180)
    ax.bar(x - width / 2, twm, width, label="TWM", color="#247C7A")
    ax.bar(x + width / 2, flus, width, label="FLUS", color="#B56A2A")
    ax.set_ylim(0, 1.0)
    ax.set_ylabel("mean score")
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False, loc="upper left")
    for i, (a, b) in enumerate(zip(twm, flus, strict=False)):
        ax.text(i - width / 2, a + 0.015, f"{a:.3f}", ha="center", va="bottom", fontsize=8)
        ax.text(i + width / 2, b + 0.015, f"{b:.3f}", ha="center", va="bottom", fontsize=8)
    ax.spines[["top", "right", "left"]].set_visible(False)
    fig.tight_layout()
    path = ASSET_DIR / "twm_flus_100case_metric_chart.png"
    fig.savefig(path, transparent=False, facecolor="white")
    plt.close(fig)
    return path


def make_change_counts_chart():
    labels = ["Hits", "False alarms", "Misses"]
    pair = np.array([80548, 134024, 167559])
    contrast = np.array([80738, 133834, 167369])
    x = np.arange(len(labels))
    width = 0.36
    fig, ax = plt.subplots(figsize=(8.8, 3.8), dpi=180)
    ax.bar(x - width / 2, pair, width, label="Pair support", color="#7A8EA4")
    ax.bar(x + width / 2, contrast, width, label="Contrast guard", color="#247C7A")
    ax.set_ylabel("100-case count")
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False, loc="upper left")
    for i, (a, b) in enumerate(zip(pair, contrast, strict=False)):
        ax.text(i - width / 2, a + 2600, f"{a:,}", ha="center", va="bottom", fontsize=8)
        ax.text(i + width / 2, b + 2600, f"{b:,}", ha="center", va="bottom", fontsize=8)
    ax.spines[["top", "right", "left"]].set_visible(False)
    fig.tight_layout()
    path = ASSET_DIR / "twm_contrast_change_counts.png"
    fig.savefig(path, transparent=False, facecolor="white")
    plt.close(fig)
    return path


def make_world_model_route_chart():
    routes = [
        "RL control",
        "Video generative",
        "Game simulation",
        "Autonomous driving",
        "Scientific model",
        "Geospatial",
    ]
    maturity = np.array([0.78, 0.66, 0.62, 0.70, 0.74, 0.40])
    deploy = np.array([0.55, 0.35, 0.48, 0.55, 0.62, 0.32])
    x = np.arange(len(routes))
    fig, ax = plt.subplots(figsize=(9.5, 4.0), dpi=180)
    ax.plot(x, maturity, marker="o", lw=2.5, color="#1B5C8C", label="model maturity")
    ax.plot(x, deploy, marker="o", lw=2.5, color="#B56A2A", label="production usability")
    ax.fill_between(x, maturity, deploy, color="#D7E0E8", alpha=0.45)
    ax.set_ylim(0, 1)
    ax.set_xticks(x)
    ax.set_xticklabels(routes, rotation=12, ha="right")
    ax.set_yticks([0, 0.25, 0.5, 0.75, 1.0])
    ax.grid(axis="y", alpha=0.22)
    ax.legend(frameon=False, loc="upper right")
    ax.spines[["top", "right", "left"]].set_visible(False)
    fig.tight_layout()
    path = ASSET_DIR / "world_model_routes_maturity.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def section_slide(prs: Presentation, idx: int, section: str, title: str, subtitle: str, color=COLORS["dark"]):
    slide = blank(prs)
    fill_background(slide, color)
    add_chip(slide, 0.72, 0.64, section, COLORS["teal"], w=1.25)
    add_text(slide, 0.72, 1.42, 10.8, 0.75, title, 31, True, COLORS["white"])
    add_text(slide, 0.76, 2.42, 9.4, 0.72, subtitle, 16, False, rgb("D7E0E8"), line_spacing=1.0)
    add_text(slide, 12.0, 6.95, 0.7, 0.24, f"{idx:02d}", 9, True, rgb("D7E0E8"), PP_ALIGN.RIGHT)
    return slide


def build_deck():
    ensure_assets()
    metric_chart = make_metric_chart()
    counts_chart = make_change_counts_chart()
    route_chart = make_world_model_route_chart()

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    slides = []

    # 1
    slide = blank(prs)
    slide.shapes.add_picture(str(ASSET_DIR / "cover_data_for_ai.png"), 0, 0, width=Inches(W), height=Inches(H))
    add_chip(slide, 0.72, 0.64, "内部技术交流", COLORS["teal"], w=1.55)
    add_text(slide, 0.72, 1.35, 8.5, 0.9, "Data for AI：\n从时空数据治理到地理空间世界模型", 28, True, COLORS["white"])
    add_text(slide, 0.76, 3.08, 7.6, 0.55, "MMFE 作为 AI-ready 时空数据底座，TWM 作为自然资源行业 Geospatial World Model 实例。", 15, False, rgb("E6EEF5"))
    add_text(slide, 0.76, 6.82, 4.8, 0.24, "GIS Data Agent / 2026-07-02", 8, False, rgb("D7E0E8"))
    add_notes(slide, "开场强调：这不是产品宣传页，而是一个技术路线分享：Data for AI 如何从治理走向模拟与规划。")
    slides.append("标题页")

    # 2
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "今天只讲一条主线", "执行摘要", 2)
    add_card(slide, 0.75, 1.35, 3.65, 3.6, "1. Data for AI", "数据治理的目标从“可查、可看”升级为“可被模型稳定消费”。\n\nMMFE 提供多源时空数据接入、融合、质量和语义治理。", COLORS["panel"], COLORS["blue"])
    add_card(slide, 4.85, 1.35, 3.65, 3.6, "2. World Model", "世界模型不是热词，而是对状态、演化规律和行动后果的可计算表达。\n\n关键能力是预测、模拟、规划。", COLORS["panel_green"], COLORS["green"])
    add_card(slide, 8.95, 1.35, 3.65, 3.6, "3. Geospatial WM", "地理空间世界模型把 GIS 对象、时序变化、空间约束和行业行动组织进一个模拟闭环。\n\nTWM 是自然资源行业实例。", COLORS["panel_amber"], COLORS["amber"])
    add_text(slide, 1.0, 5.55, 11.5, 0.5, "一句话：MMFE 解决“AI 可用的数据”，TWM 解决“数据驱动的世界状态推演”。", 20, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "先把听众拉到同一语境：Data for AI 不只是建湖仓，而是为后续世界模型提供观测、状态和证据。")
    slides.append("主线")

    # 3
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "为什么传统 GIS 数据治理不够了", "Data for AI 的问题背景", 3)
    problems = [
        ("面向人看", "传统治理面向制图、检索、统计报表；AI 需要稳定 API、结构化语义和可追溯 lineage。"),
        ("面向单数据集", "AI 场景通常同时需要矢量、栅格、表格、影像、规则、审批历史和文本证据。"),
        ("面向静态状态", "世界模型要求看到历史变化、约束条件和行动干预，而不只是当前底图。"),
        ("缺少模型反馈", "数据质量、融合冲突、语义错配需要进入模型和 agent 的反馈闭环。"),
    ]
    for i, (t, b) in enumerate(problems):
        x = 0.75 + (i % 2) * 6.15
        y = 1.35 + (i // 2) * 2.15
        add_card(slide, x, y, 5.55, 1.55, t, b, COLORS["white"], [COLORS["blue"], COLORS["teal"], COLORS["amber"], COLORS["red"]][i], 15, 11)
    add_text(slide, 0.82, 6.15, 11.8, 0.36, "AI 时代的数据治理对象变成：数据 + 语义 + 质量 + 证据 + 时间演化 + 可行动接口。", 17, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("传统治理缺口")

    # 4
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "Data for AI 的分层目标", "从数据资产到模型燃料", 4)
    stages = [
        ("Data", "多源接入\n遥感/矢量/栅格/表格"),
        ("Governance", "标准化\n质量/血缘/权限"),
        ("Semantics", "实体/关系/本体\n字段语义对齐"),
        ("Model-ready", "特征/样本/时序\n稳定 contract"),
        ("World Model", "预测/模拟/规划\n证据门控"),
    ]
    for i, (t, b) in enumerate(stages):
        x = 0.75 + i * 2.45
        add_panel(slide, x, 2.0, 1.9, 1.55, COLORS["white"])
        add_text(slide, x + 0.12, 2.22, 1.66, 0.28, t, 14, True, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i], PP_ALIGN.CENTER)
        add_text(slide, x + 0.12, 2.72, 1.66, 0.48, b, 9, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_arrow(slide, x + 1.95, 2.78, x + 2.35, 2.78, COLORS["line"], 1.4)
    add_card(slide, 1.0, 4.55, 5.2, 1.15, "MMFE 的位置", "把多源时空数据治理成 AI-ready 数据产品。", COLORS["panel"], COLORS["blue"], 14, 12)
    add_card(slide, 7.1, 4.55, 5.2, 1.15, "TWM 的位置", "在 AI-ready 数据之上学习状态演化和行动后果。", COLORS["panel_green"], COLORS["green"], 14, 12)
    add_footer(slide)
    slides.append("Data for AI 分层")

    # 5
    slide = section_slide(prs, 5, "PART 1", "世界模型：先把概念说清楚", "不同团队对 world model 的理解不一样，部门内部交流必须先建立共同词典。", COLORS["dark"])
    add_notes(slide, "过渡：先讲世界模型，再讲地理空间世界模型，最后落到 TWM。")
    slides.append("章节：世界模型")

    # 6
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "世界模型不是一个单一产品名", "概念定义", 6)
    add_text(slide, 0.85, 1.25, 11.7, 0.72, "世界模型 = 对现实环境的状态、演化规律、行动影响的可计算表达", 26, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_panel(slide, 1.0, 2.45, 11.3, 2.75, COLORS["white"])
    triples = [
        ("State", "世界现在是什么状态", "对象、关系、属性、约束、证据"),
        ("Dynamics", "世界会如何变化", "转移、趋势、物理/空间过程、不确定性"),
        ("Action", "行动会造成什么后果", "反事实、规划动作、收益、风险"),
    ]
    for i, (t, q, b) in enumerate(triples):
        x = 1.45 + i * 3.65
        add_chip(slide, x, 2.85, t, [COLORS["blue"], COLORS["green"], COLORS["amber"]][i], w=1.1)
        add_text(slide, x, 3.35, 2.55, 0.3, q, 14, True, COLORS["ink"], PP_ALIGN.CENTER)
        add_text(slide, x, 3.88, 2.55, 0.52, b, 10, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(slide, x + 2.75, 3.65, x + 3.25, 3.65, COLORS["line"], 1.2)
    add_text(slide, 1.05, 5.8, 11.1, 0.42, "它和 LLM 的区别：LLM 主要预测文本序列；世界模型要预测环境状态变化，并支持模拟与规划。", 15, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("世界模型定义")

    # 7
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "世界模型的主要技术路线", "路线版图", 7)
    headers = ["路线", "代表输入", "核心目标", "成熟度判断"]
    xs = [0.75, 2.75, 5.25, 8.35]
    ws = [1.7, 2.15, 2.8, 4.1]
    for x, w, h in zip(xs, ws, headers, strict=False):
        add_panel(slide, x, 1.15, w, 0.45, COLORS["dark"], radius=False)
        add_text(slide, x + 0.06, 1.27, w - 0.12, 0.14, h, 8, True, COLORS["white"], PP_ALIGN.CENTER)
    rows = [
        ("RL / 控制", "状态 + 动作", "想象未来并学习策略", "Dreamer 等路线较成熟，仍依赖任务环境定义"),
        ("视频生成", "视频/图像/文本", "生成未来场景", "表现力强，可控性和物理一致性仍是瓶颈"),
        ("游戏/交互", "视频 + 动作", "生成可交互环境", "Genie/GameNGen 推动明显，真实业务迁移仍早期"),
        ("自动驾驶", "相机/LiDAR/动作", "驾驶场景预测与仿真", "行业投入大，但安全验证门槛极高"),
        ("科学/物理", "观测 + 方程", "物理系统模拟", "气象、材料、地球系统进展快，依赖领域约束"),
        ("Geospatial", "GIS + 遥感 + 规则", "空间状态演化与行动后果", "强 GIS 行业需求明确，工程范式仍在形成"),
    ]
    for i, row in enumerate(rows):
        y = 1.68 + i * 0.72
        fill = COLORS["white"] if i % 2 == 0 else rgb("F9FBFC")
        for x, w, text in zip(xs, ws, row, strict=False):
            add_panel(slide, x, y, w, 0.58, fill, radius=False)
            add_text(slide, x + 0.08, y + 0.13, w - 0.14, 0.2, text, 7, False, COLORS["ink"] if x == xs[0] else COLORS["muted"])
    add_footer(slide)
    slides.append("世界模型路线")

    # 8
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "各路线的发展状态：热，但还没统一范式", "路线成熟度", 8)
    add_image_fit(slide, route_chart, 0.8, 1.25, 6.7, 3.25)
    add_card(slide, 8.05, 1.25, 4.45, 1.2, "共同趋势", "从感知模型转向可预测、可交互、可规划的环境模型。", COLORS["panel"], COLORS["blue"], 14, 11)
    add_card(slide, 8.05, 2.85, 4.45, 1.2, "共同瓶颈", "长时一致性、物理/规则约束、可验证性、真实世界反馈闭环。", COLORS["panel_amber"], COLORS["amber"], 14, 11)
    add_card(slide, 8.05, 4.45, 4.45, 1.2, "GWM 机会", "地理空间任务天然有对象、边界、规则、时间序列和规划行动。", COLORS["panel_green"], COLORS["green"], 14, 11)
    add_text(slide, 0.9, 5.7, 6.4, 0.38, "注：图中 maturity/usability 为技术路线判断示意，不是外部 benchmark 评分。", 8, False, COLORS["muted"])
    add_footer(slide)
    slides.append("世界模型发展状态")

    # 9
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "为什么 GIS 行业需要 Geospatial World Model", "从描述世界到推演世界", 9)
    left = [
        ("传统 GIS", "回答“哪里有什么”\n图层、查询、叠加、统计、制图"),
        ("数字孪生", "回答“现在长什么样”\n可视化、监控、状态同步"),
        ("Geospatial WM", "回答“会怎么变、行动会怎样”\n预测、模拟、反事实、规划"),
    ]
    for i, (t, b) in enumerate(left):
        add_card(slide, 0.95 + i * 4.05, 1.55, 3.35, 2.05, t, b, [COLORS["panel"], COLORS["panel_amber"], COLORS["panel_green"]][i], [COLORS["blue"], COLORS["amber"], COLORS["green"]][i], 16, 12)
        if i < 2:
            add_arrow(slide, 4.15 + i * 4.05, 2.58, 4.62 + i * 4.05, 2.58, COLORS["line"], 1.5)
    add_panel(slide, 1.0, 4.55, 11.3, 1.2, COLORS["white"])
    add_text(slide, 1.22, 4.82, 10.85, 0.36, "强 GIS 行业的共性：空间对象 + 时序变化 + 业务规则 + 行动干预 + 规划目标。", 18, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("为什么 GWM")

    # 10
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "Geospatial World Model 的定义", "面向地理空间状态演化与行动后果", 10)
    add_text(slide, 0.9, 1.25, 11.6, 0.86, "Geospatial World Model 是面向地理空间对象、时空状态演化和空间行动后果的世界模型。", 24, True, COLORS["blue"], PP_ALIGN.CENTER)
    items = [
        ("对象", "地块、道路、河流、设施、项目、行政区"),
        ("关系", "邻接、包含、重叠、连通、上下游、服务范围"),
        ("过程", "土地变化、城市扩张、灾害扩散、交通拥堵"),
        ("行动", "保护、开发、审批、调度、建设、管控"),
        ("证据", "数据来源、规则命中、质量、血缘、适用边界"),
    ]
    for i, (t, b) in enumerate(items):
        x = 0.82 + i * 2.48
        add_panel(slide, x, 2.8, 2.0, 1.75, COLORS["white"])
        add_chip(slide, x + 0.48, 3.05, t, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i], w=0.82)
        add_text(slide, x + 0.18, 3.62, 1.64, 0.42, b, 9, False, COLORS["muted"], PP_ALIGN.CENTER)
    add_text(slide, 0.95, 5.68, 11.35, 0.46, "GWM 的目标不是替代 GIS，而是让 GIS 数据进入预测、模拟、规划和审计闭环。", 17, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("GWM 定义")

    # 11
    slide = section_slide(prs, 11, "PART 2", "MMFE：时空数据治理作为 AI-ready 底座", "没有稳定、语义化、可追溯的数据产品，世界模型只能停留在 demo。", COLORS["blue"])
    slides.append("章节：MMFE")

    # 12
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "MMFE 在 GIS Data Agent 中的定位", "多模态时空数据治理底座", 12)
    add_image_fit(slide, ASSET_DIR / "twm_upstream_foundation_overview.png", 0.7, 1.08, 5.7, 4.35)
    add_card(slide, 6.85, 1.18, 5.55, 0.98, "定位", "把矢量、栅格、表格、遥感、业务数据治理成 AI 可用的数据产品。", COLORS["panel"], COLORS["blue"], 14, 11)
    add_card(slide, 6.85, 2.48, 5.55, 0.98, "关键能力", "数据探测、兼容性评估、语义对齐、融合执行、质量验证。", COLORS["panel_green"], COLORS["green"], 14, 11)
    add_card(slide, 6.85, 3.78, 5.55, 0.98, "对 TWM 的价值", "提供状态观测、训练样本、空间对象、规则证据和质量边界。", COLORS["panel_amber"], COLORS["amber"], 14, 11)
    add_text(slide, 0.85, 5.82, 11.8, 0.32, "MMFE 是 Data for AI 的生产层：先把“数据能用”变成“模型敢用”。", 17, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("MMFE 定位")

    # 13
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "MMFE 五阶段流水线", "从多源数据到可信融合结果", 13)
    steps = [
        ("1 数据探测", "CRS / 范围 / 字段 / 分辨率"),
        ("2 兼容性评估", "空间重叠 / 坐标冲突 / 字段匹配"),
        ("3 语义对齐", "本体 / 等价组 / LLM schema / 单位"),
        ("4 融合执行", "join / overlay / zonal / pushdown"),
        ("5 质量验证", "空值 / 拓扑 / 分布偏移 / 冲突"),
    ]
    for i, (t, b) in enumerate(steps):
        x = 0.75 + i * 2.45
        add_panel(slide, x, 1.55, 1.95, 2.05, COLORS["white"])
        add_chip(slide, x + 0.2, 1.9, t, [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i], w=1.55)
        add_text(slide, x + 0.17, 2.56, 1.6, 0.45, b, 9, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < 4:
            add_arrow(slide, x + 1.98, 2.58, x + 2.34, 2.58, COLORS["line"], 1.2)
    add_panel(slide, 1.0, 4.55, 11.3, 1.1, COLORS["panel"])
    add_text(slide, 1.25, 4.88, 10.85, 0.28, "关键变化：数据治理从“手工整理”变成可被 agent 调用、可复核、可回放的流水线。", 16, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("MMFE 五阶段")

    # 14
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "MMFE 不只是融合算子，而是治理 contract", "AI-ready 数据产品", 14)
    add_card(slide, 0.75, 1.2, 3.85, 1.45, "语义 contract", "字段、单位、业务概念、空间对象和本体关系要稳定可解释。", COLORS["panel"], COLORS["blue"], 14, 11)
    add_card(slide, 4.95, 1.2, 3.85, 1.45, "质量 contract", "质量门控、冲突记录、置信度、异常值和拓扑问题要结构化。", COLORS["panel_green"], COLORS["green"], 14, 11)
    add_card(slide, 9.15, 1.2, 3.35, 1.45, "血缘 contract", "来源、参数、版本、处理链路和适用边界要可追踪。", COLORS["panel_amber"], COLORS["amber"], 14, 11)
    add_panel(slide, 1.15, 3.55, 10.95, 1.45, COLORS["white"])
    add_text(slide, 1.35, 3.86, 10.55, 0.38, "World Model 消费的不是“文件”，而是带语义、质量和血缘的状态观测。", 20, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("MMFE contract")

    # 15
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "MMFE 展示效果：融合质量进入可解释界面", "项目截图", 15)
    if (ASSET_DIR / "mmfe_fusion_quality_e2e.png").exists():
        add_image_fit(slide, ASSET_DIR / "mmfe_fusion_quality_e2e.png", 0.72, 1.05, 7.15, 5.35)
    add_card(slide, 8.28, 1.25, 4.15, 1.25, "从结果到证据", "不是只输出 GeoJSON，而是把质量、冲突、来源和解释交给用户与 agent。", COLORS["panel"], COLORS["blue"], 14, 11)
    add_card(slide, 8.28, 2.95, 4.15, 1.25, "从工具到闭环", "质量问题可以进入后续 schema 对齐、融合策略和数据修复。", COLORS["panel_green"], COLORS["green"], 14, 11)
    add_card(slide, 8.28, 4.65, 4.15, 1.25, "对 TWM 的意义", "TWM 的训练样本和状态观测必须携带这种可验证边界。", COLORS["panel_amber"], COLORS["amber"], 14, 11)
    add_footer(slide)
    slides.append("MMFE 截图")

    # 16
    slide = section_slide(prs, 16, "PART 3", "TWM：自然资源行业的 Geospatial World Model 实例", "TWM 不是泛化概念包装，而是把自然资源治理问题具体化为可运行的 GWM 架构。", COLORS["green"])
    slides.append("章节：TWM")

    # 17
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "TWM 的行业定义", "Geospatial World Model in Natural Resources", 17)
    add_text(slide, 0.9, 1.18, 11.4, 0.68, "TWM 是 Geospatial World Model 在自然资源行业中的一个实例。", 26, True, COLORS["green"], PP_ALIGN.CENTER)
    cards = [
        ("建模对象", "国土空间状态、土地覆盖/利用、地块、项目、规则和证据。"),
        ("建模过程", "土地变化、规划管控、保护/开发行动、约束风险和效用变化。"),
        ("输出能力", "未来状态、变化位置、约束风险、方案比较和审计边界。"),
    ]
    for i, (t, b) in enumerate(cards):
        add_card(slide, 0.9 + i * 4.1, 2.55, 3.55, 1.75, t, b, [COLORS["panel"], COLORS["panel_green"], COLORS["panel_amber"]][i], [COLORS["blue"], COLORS["green"], COLORS["amber"]][i], 15, 11)
    add_text(slide, 1.05, 5.45, 11.2, 0.46, "它可以扩展到其它强 GIS 行业，但自然资源是当前最适合落地的第一行业实例。", 17, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("TWM 定义")

    # 18
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "TWM 三层能力：Renderer / Simulator / Planner", "不要把三者混为一谈", 18)
    add_card(slide, 0.85, 1.3, 3.65, 3.25, "Renderer", "把 GIS 状态、规则、证据、风险和模拟结果渲染成可观察、可审计的业务视图。\n\n不负责核心预测。", COLORS["panel"], COLORS["blue"], 15, 11)
    add_card(slide, 4.85, 1.3, 3.65, 3.25, "Simulator", "TWM 核心。\n从历史状态学习变化规律，预测未来状态、约束风险和不确定性。\n\n与 FLUS 对比的是这部分。", COLORS["panel_green"], COLORS["green"], 15, 11)
    add_card(slide, 8.85, 1.3, 3.65, 3.25, "Planner", "消费 simulator 输出，进行方案搜索、约束检查、候选排序和复核任务生成。\n\n不是本次 FLUS 对比主体。", COLORS["panel_amber"], COLORS["amber"], 15, 11)
    add_text(slide, 0.95, 5.6, 11.5, 0.35, "当前最强创新点在 Simulator：action / evidence / topology / demand 共同约束空间状态推演。", 17, True, COLORS["green"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("TWM 三层")

    # 19
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "TWM Simulator 的算法链路", "与 GeoSOS-FLUS 对比的核心能力", 19)
    chain = [
        ("历史状态", "train_start\ntrain_end"),
        ("Transition scoring", "source -> target\n变化倾向"),
        ("Replay calibration", "训练期回放\n识别误报 pair"),
        ("Topology guards", "稳定区/前沿/邻域\n空间支撑"),
        ("Demand allocation", "需求守恒\n变化预算"),
        ("Forecast map", "holdout 评估\n指标聚合"),
    ]
    for i, (t, b) in enumerate(chain):
        x = 0.6 + i * 2.1
        add_panel(slide, x, 2.0, 1.68, 1.4, COLORS["white"])
        add_text(slide, x + 0.1, 2.22, 1.48, 0.26, t, 9, True, COLORS["ink"], PP_ALIGN.CENTER)
        add_text(slide, x + 0.1, 2.72, 1.48, 0.42, b, 8, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < len(chain) - 1:
            add_arrow(slide, x + 1.72, 2.68, x + 2.02, 2.68, COLORS["line"], 1.2)
    add_card(slide, 1.15, 4.55, 5.3, 1.2, "最新关键机制", "pair topology support contrast：高误报 source->target pair 中，unsupported cells 降分，supported cells 加分。", COLORS["panel_green"], COLORS["green"], 14, 11)
    add_card(slide, 7.0, 4.55, 5.3, 1.2, "严格边界", "训练只使用 train_start -> train_end；holdout 只用于 evaluation，不进入训练。", COLORS["panel_amber"], COLORS["amber"], 14, 11)
    add_footer(slide)
    slides.append("TWM 算法链路")

    # 20
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "TWM 与 GeoSOS-FLUS：不是同一种模型包装", "比较边界", 20)
    add_card(slide, 0.85, 1.25, 5.55, 3.9, "GeoSOS-FLUS", "经典土地利用变化模拟路线：\n\n- ANN suitability\n- Cellular automata allocation\n- adaptive inertia\n- land-use competition\n- scenario demand\n\n优势：整体地图保持更稳，成熟度高。", COLORS["panel_amber"], COLORS["amber"], 15, 11)
    add_card(slide, 6.95, 1.25, 5.55, 3.9, "TWM Simulator", "地理空间世界模型路线：\n\n- train-only transition replay\n- topology stability/support guards\n- pair-specific contrast ranking\n- demand-constrained allocation\n- evidence-gated comparison\n\n优势：变化定位能力更强。", COLORS["panel_green"], COLORS["green"], 15, 11)
    add_text(slide, 0.95, 5.78, 11.5, 0.38, "严谨说法：当前 TWM 在变化指标上超过固定 FLUS baseline，不是全面替代 FLUS。", 17, True, COLORS["red"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("TWM vs FLUS 边界")

    # 21
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "对比是否公平：同一数据、同一切分、同一评估", "实验协议", 21)
    protocol = [
        ("数据", "Dynamic World annual 100m raster cases"),
        ("时间切分", "train_start -> train_end 训练；train_end -> holdout 评估"),
        ("输入地图", "同一个 initial map、holdout map、evaluation mask"),
        ("隔离", "holdout 只用于 evaluation，不进入 TWM training"),
        ("复用", "后续 TWM 迭代复用固定 FLUS 输出，避免随机性"),
    ]
    for i, (t, b) in enumerate(protocol):
        y = 1.18 + i * 0.88
        add_panel(slide, 1.0, y, 2.1, 0.55, COLORS["dark"], radius=False)
        add_text(slide, 1.08, y + 0.16, 1.9, 0.14, t, 8, True, COLORS["white"], PP_ALIGN.CENTER)
        add_panel(slide, 3.1, y, 8.95, 0.55, COLORS["white"], radius=False)
        add_text(slide, 3.32, y + 0.16, 8.45, 0.14, b, 9, False, COLORS["ink"])
    add_text(slide, 1.1, 6.0, 10.9, 0.38, "公平不等于内部机制完全相同；公平在于数据来源、时间切分、评估目标和 holdout 隔离一致。", 15, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("公平协议")

    # 22
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "100-case 对比结果：TWM 强在变化定位", "TWM vs FLUS 指标", 22)
    add_image_fit(slide, metric_chart, 0.75, 1.15, 7.5, 4.65)
    add_card(slide, 8.65, 1.32, 3.85, 1.2, "TWM 赢的指标", "Change FoM +0.044707\nChange F1 +0.070240", COLORS["panel_green"], COLORS["green"], 14, 12)
    add_card(slide, 8.65, 2.92, 3.85, 1.2, "TWM 仍落后的指标", "OA -0.017399\nKappa -0.039776\nMacro-F1 -0.020851", COLORS["panel_red"], COLORS["red"], 14, 12)
    add_card(slide, 8.65, 4.52, 3.85, 1.2, "结论边界", "变化预测能力领先；整体地图模拟能力尚未全面领先。", COLORS["panel_amber"], COLORS["amber"], 14, 11)
    add_footer(slide)
    slides.append("对比指标")

    # 23
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "为什么 pair topology support contrast 有效", "固定变化配额下改进空间排序", 23)
    add_image_fit(slide, counts_chart, 0.8, 1.15, 6.8, 3.65)
    add_card(slide, 8.1, 1.18, 4.35, 1.18, "不是增加变化量", "Predicted changes 保持 214,572 不变。", COLORS["panel"], COLORS["blue"], 14, 12)
    add_card(slide, 8.1, 2.72, 4.35, 1.18, "而是排序更准", "hits +190，false alarms -190，misses -190。", COLORS["panel_green"], COLORS["green"], 14, 12)
    add_card(slide, 8.1, 4.26, 4.35, 1.18, "地理空间特性", "对有前沿/邻域支撑的变化加分，对 unsupported cells 降分。", COLORS["panel_amber"], COLORS["amber"], 14, 11)
    add_text(slide, 0.95, 5.75, 11.3, 0.36, "这体现的是 simulator 的 geospatial topology ranking 能力，不是 renderer 或 planner 的能力。", 16, True, COLORS["green"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("contrast guard")

    # 24
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "从 TWM 到其它强 GIS 行业", "Geospatial World Model 可迁移方向", 24)
    industries = [
        ("自然资源", "土地变化 / 生态保护 / 规划管控"),
        ("城市治理", "城市扩张 / 设施布局 / 人口活动"),
        ("农业", "作物生长 / 耕地变化 / 产量模拟"),
        ("水利", "洪水演进 / 流域调度 / 风险区变化"),
        ("应急", "灾害扩散 / 避难路径 / 资源调度"),
        ("交通", "路网状态 / 拥堵演化 / 建设影响"),
        ("能源", "选址 / 负荷分布 / 管网风险"),
        ("生态环境", "污染扩散 / 生态修复 / 碳汇评估"),
    ]
    for i, (t, b) in enumerate(industries):
        x = 0.78 + (i % 4) * 3.1
        y = 1.25 + (i // 4) * 2.05
        add_card(slide, x, y, 2.65, 1.45, t, b, COLORS["white"], [COLORS["blue"], COLORS["teal"], COLORS["green"], COLORS["amber"], COLORS["red"], COLORS["purple"], COLORS["blue"], COLORS["green"]][i], 13, 9)
    add_text(slide, 0.95, 5.95, 11.4, 0.34, "TWM 是自然资源行业实例；Geospatial World Model 是可跨行业复用的技术范式。", 16, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("行业扩展")

    # 25
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "MMFE 与 TWM 的端到端关系", "系统总架构", 25)
    layers = [
        ("Data Sources", "遥感 / 矢量 / 栅格 / 表格 / 业务规则"),
        ("MMFE", "探测 / 对齐 / 融合 / 质量 / 血缘"),
        ("Semantic Products", "空间实体 / 时序样本 / 规则证据 / 特征表"),
        ("TWM Simulator", "状态转移 / topology guard / demand allocation"),
        ("Planner & Apps", "方案比较 / 风险解释 / 审计 / 可视化"),
    ]
    for i, (t, b) in enumerate(layers):
        x = 0.7 + i * 2.48
        add_panel(slide, x, 2.0, 1.95, 1.65, [COLORS["panel"], COLORS["panel_green"], COLORS["white"], COLORS["panel_amber"], COLORS["panel"]][i])
        add_text(slide, x + 0.1, 2.25, 1.75, 0.24, t, 10, True, COLORS["ink"], PP_ALIGN.CENTER)
        add_text(slide, x + 0.14, 2.78, 1.65, 0.45, b, 8, False, COLORS["muted"], PP_ALIGN.CENTER)
        if i < 4:
            add_arrow(slide, x + 1.98, 2.83, x + 2.36, 2.83, COLORS["line"], 1.2)
    add_panel(slide, 1.05, 4.65, 11.2, 0.95, COLORS["dark"])
    add_text(slide, 1.3, 4.95, 10.75, 0.3, "Data for AI 的落点：数据治理不是终点，模型可消费、可验证、可规划才是闭环。", 17, True, COLORS["white"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("端到端架构")

    # 26
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "部门内部可以落地的价值", "从平台能力到业务工作流", 26)
    value = [
        ("数据资产管理升级", "从文件/图层治理升级为 AI-ready 数据产品治理。"),
        ("模型验证更严格", "训练、评估、证据边界和适用范围进入统一报告。"),
        ("业务推演能力增强", "从查图查数走向变化预测、方案模拟和风险解释。"),
        ("Agent 工作流闭环", "自然语言入口连接数据治理、模型模拟、规划输出和审计。"),
    ]
    for i, (t, b) in enumerate(value):
        x = 0.85 + (i % 2) * 6.05
        y = 1.35 + (i // 2) * 2.0
        add_card(slide, x, y, 5.35, 1.35, t, b, COLORS["white"], [COLORS["blue"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i], 15, 11)
    add_text(slide, 0.95, 5.7, 11.4, 0.36, "建议内部试点：选一个稳定业务场景，先跑 MMFE 数据 contract，再接 TWM simulator 评估闭环。", 16, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_footer(slide)
    slides.append("部门价值")

    # 27
    slide = blank(prs)
    fill_background(slide)
    add_title(slide, "后续技术路线", "从工程可用到模型可信", 27)
    roadmap = [
        ("P1 数据底座", "生产级 MMFE lakehouse、语义层、质量与血缘治理"),
        ("P2 模拟器增强", "region/class-conditioned topology support，学习化替代固定阈值"),
        ("P3 规划闭环", "planner 消费 simulator 输出，形成方案比较和复核任务"),
        ("P4 行业扩展", "从自然资源 TWM 扩展到城市、水利、应急、农业等 GWM"),
    ]
    for i, (t, b) in enumerate(roadmap):
        y = 1.24 + i * 1.2
        add_chip(slide, 0.92, y + 0.12, f"{i + 1}", [COLORS["blue"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i], w=0.42)
        accent = [COLORS["blue"], COLORS["green"], COLORS["amber"], COLORS["purple"]][i]
        add_panel(slide, 1.55, y, 10.8, 0.92, COLORS["white"])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.55), Inches(y), Inches(0.08), Inches(0.92))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        add_text(slide, 1.82, y + 0.18, 10.25, 0.22, t, 13, True, COLORS["ink"])
        add_text(slide, 1.82, y + 0.56, 10.25, 0.2, b, 9, False, COLORS["muted"])
    add_footer(slide)
    slides.append("后续路线")

    # 28
    slide = blank(prs)
    fill_background(slide, COLORS["dark"])
    add_text(slide, 0.85, 0.82, 10.8, 0.62, "三个 takeaway", 30, True, COLORS["white"])
    takeaway = [
        ("1", "Data for AI 不是数据平台换名，而是让数据成为模型可消费、可验证的状态观测。"),
        ("2", "世界模型的关键不是生成画面，而是学习状态、动态和行动后果。"),
        ("3", "TWM 是自然资源行业 GWM 实例；当前优势在 simulator 的变化预测能力。"),
    ]
    for i, (n, text) in enumerate(takeaway):
        y = 2.0 + i * 1.15
        add_chip(slide, 1.0, y, n, COLORS["teal"], w=0.45)
        add_text(slide, 1.72, y + 0.05, 10.1, 0.32, text, 16, True, rgb("E6EEF5"))
    add_text(slide, 0.9, 6.55, 11.6, 0.3, "讨论：如何选择一个部门场景，把 MMFE 数据治理和 TWM 模拟器接成可复核闭环？", 15, True, COLORS["white"], PP_ALIGN.CENTER)
    add_footer(slide, "References: Ha & Schmidhuber 2018; DreamerV3 2023/2024; Genie 2024; GameNGen 2024; GAIA-1 2023; AlphaEarth Foundations 2025; GIS Data Agent local docs.")
    slides.append("总结")

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if not run.font.name:
                            run.font.name = FONT_CN

    prs.save(PPTX_PATH)
    OUTLINE_PATH.write_text("\n".join(f"{i + 1}. {name}" for i, name in enumerate(slides)) + "\n", encoding="utf-8")
    QA_PATH.write_text(
        "\n".join(
            [
                "# QA Report",
                "",
                f"- PPTX: `{PPTX_PATH.name}`",
                f"- Slide count: {len(slides)}",
                "- Generation: python-pptx, project screenshots, generated charts.",
                "- Visual QA to run: LibreOffice PDF export plus rendered page contact sheets.",
                "",
                "## Key Content Checks",
                "",
                "- Explains world model definition and route landscape.",
                "- Defines Geospatial World Model and positions TWM as the natural-resource instance.",
                "- Uses MMFE as the AI-ready spatiotemporal data-governance foundation.",
                "- States TWM/FLUS boundary: change metrics lead, OA/Kappa/Macro-F1 still trail.",
                "- Distinguishes renderer, simulator and planner; FLUS comparison targets simulator.",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    return PPTX_PATH


if __name__ == "__main__":
    path = build_deck()
    print(path)
