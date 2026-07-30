#!/usr/bin/env python3
"""Build the evidence-backed Gemma 4 finals deck and its visual assets."""

from __future__ import annotations

import argparse
import csv
import json
from pathlib import Path
from typing import Any

import geopandas as gpd
import matplotlib.pyplot as plt
from matplotlib import font_manager
from matplotlib.patches import Patch
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FINALS_DIR = ROOT / "docs" / "finals"
ASSET_DIR = FINALS_DIR / "assets"
DEFAULT_TEMPLATE = Path.home() / "Downloads" / "路演框架（GDG Shanghai）.pptx"
DEFAULT_OUTPUT = FINALS_DIR / "GIS_Data_Agent_Gemma4_Finals_CN.pptx"

BENCHMARK_CSV = ROOT / "docs" / "assets" / "gemma4_host228_scale_sweep_summary.csv"
RELIABILITY_JSON = (
    ROOT
    / "data_agent"
    / "demo_evidence"
    / "paper9"
    / "finals_20260730"
    / "adk_reliability_report.json"
)
PREFLIGHT_JSON = RELIABILITY_JSON.with_name("finals_preflight_final.json")
FARMLAND_BASE_RUN_DIR = (
    ROOT
    / "data_agent"
    / "uploads"
    / "finals_verification_fixed"
    / "world_model_v21"
    / "20260730_003511_666173"
)
FARMLAND_EVIDENCE_DIR = (
    FINALS_DIR / "evidence" / "world_model_bishan_20260730_155442"
)
FARMLAND_SUMMARY_JSON = FARMLAND_EVIDENCE_DIR / "mpc_summary.json"
FARMLAND_AUDIT_JSON = FARMLAND_EVIDENCE_DIR / "paper9_agent_audit.json"
FARMLAND_CHANGES_GEOJSON = FARMLAND_EVIDENCE_DIR / "optimized_changes.geojson"
FARMLAND_BASE_SHP = FARMLAND_BASE_RUN_DIR / "optimized_dltb.shp"
NL2SQL_SCREENSHOT = (
    ROOT
    / "tests"
    / "e2e"
    / "artifacts"
    / "nl2sql_longest_bridge_map_2026_07_30"
    / "result.png"
)

FONT = "Microsoft YaHei"
WIDE = Inches(13.333333)
HIGH = Inches(7.5)

COLORS = {
    "bg": RGBColor(247, 249, 252),
    "white": RGBColor(255, 255, 255),
    "ink": RGBColor(22, 29, 44),
    "muted": RGBColor(91, 102, 121),
    "line": RGBColor(219, 226, 235),
    "navy": RGBColor(20, 31, 51),
    "blue": RGBColor(66, 133, 244),
    "green": RGBColor(52, 168, 83),
    "yellow": RGBColor(251, 188, 5),
    "red": RGBColor(234, 67, 53),
    "cyan": RGBColor(0, 137, 123),
    "pale_blue": RGBColor(232, 240, 254),
    "pale_green": RGBColor(231, 244, 234),
    "pale_yellow": RGBColor(254, 247, 224),
    "pale_red": RGBColor(253, 235, 233),
    "pale_cyan": RGBColor(226, 242, 241),
}


def color(name: str) -> RGBColor:
    return COLORS[name]


def load_json(path: Path) -> dict[str, Any]:
    with path.open(encoding="utf-8") as handle:
        return json.load(handle)


def load_benchmark() -> list[dict[str, str]]:
    with BENCHMARK_CSV.open(encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def finals_evidence() -> dict[str, Any]:
    benchmark = load_benchmark()
    reliability = load_json(RELIABILITY_JSON)
    summary = load_json(FARMLAND_SUMMARY_JSON)
    audit = load_json(FARMLAND_AUDIT_JSON)
    preflight = load_json(PREFLIGHT_JSON)
    runs = reliability["runs"]
    result = summary["results"][0]
    spatial_output = summary["shapefile_output"]
    if not audit.get("hard_constraint_passed"):
        raise ValueError("Finals farmland evidence did not pass the hard-constraint audit")
    if spatial_output["n_farm_to_forest"] != result["swaps_completed"]:
        raise ValueError("Farmland-to-forest feature count disagrees with swaps_completed")
    if spatial_output["n_forest_to_farm"] != result["swaps_completed"]:
        raise ValueError("Forest-to-farmland feature count disagrees with swaps_completed")
    return {
        "benchmark": benchmark,
        "reliability": reliability,
        "reliability_passed": sum(bool(run["passed"]) for run in runs),
        "reliability_total": len(runs),
        "summary": summary,
        "result": result,
        "audit": audit,
        "preflight": preflight,
    }


def ensure_assets(evidence: dict[str, Any]) -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    benchmark_path = ASSET_DIR / "gemma4_model_selection.png"
    map_path = ASSET_DIR / "farmland_mpc_change_map.png"
    nl2sql_path = ASSET_DIR / "nl2geosql_demo_crop.png"
    build_benchmark_chart(evidence["benchmark"], benchmark_path)
    if FARMLAND_BASE_SHP.exists():
        build_farmland_optimization_map(
            FARMLAND_BASE_SHP,
            FARMLAND_CHANGES_GEOJSON,
            evidence["summary"],
            map_path,
        )
    elif not map_path.exists():
        raise FileNotFoundError(
            "The base Shapefile is unavailable and no evidence-backed map asset exists"
        )

    if NL2SQL_SCREENSHOT.exists():
        build_nl2sql_crop(NL2SQL_SCREENSHOT, nl2sql_path)
    elif not nl2sql_path.exists():
        raise FileNotFoundError(
            "The NL2SQL screenshot is unavailable and no verified crop asset exists"
        )
    return {"benchmark": benchmark_path, "map": map_path, "nl2sql": nl2sql_path}


def chart_font() -> font_manager.FontProperties:
    candidates = [
        Path("/System/Library/Fonts/STHeiti Light.ttc"),
        Path("/System/Library/Fonts/STHeiti Medium.ttc"),
        Path("/System/Library/Fonts/Supplemental/Arial Unicode.ttf"),
    ]
    selected = next((path for path in candidates if path.exists()), None)
    return (
        font_manager.FontProperties(fname=str(selected))
        if selected
        else font_manager.FontProperties()
    )


def build_benchmark_chart(rows: list[dict[str, str]], output: Path) -> None:
    labels = [row["model"].replace("Gemma4:", "") for row in rows]
    accuracy = [float(row["full_ex_pct"]) for row in rows]
    runtime = [float(row["full_minutes"]) for row in rows]
    selected = labels.index("26b")
    font = chart_font()

    fig, ax = plt.subplots(figsize=(12.5, 5.6), dpi=160)
    fig.patch.set_facecolor("#F7F9FC")
    ax.set_facecolor("#F7F9FC")
    bars = ax.bar(
        labels,
        accuracy,
        color=["#D7DFEA" if index != selected else "#4285F4" for index in range(len(labels))],
        width=0.58,
        zorder=3,
    )
    ax.set_ylim(55, 100)
    ax.set_ylabel("CQ-125 执行正确率 (%)", fontproperties=font, color="#5B6679")
    ax.grid(axis="y", color="#E4E9F0", linewidth=0.8, zorder=0)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", length=0, colors="#5B6679")
    ax.tick_params(axis="x", length=0, colors="#161D2C", labelsize=11)
    for index, bar in enumerate(bars):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + 0.9,
            f"{accuracy[index]:.1f}%",
            ha="center",
            va="bottom",
            fontsize=11,
            color="#4285F4" if index == selected else "#5B6679",
            fontweight="bold" if index == selected else "normal",
        )

    runtime_ax = ax.twinx()
    runtime_ax.plot(
        labels, runtime, color="#EA4335", marker="o", linewidth=2.2, markersize=6, zorder=4
    )
    runtime_ax.set_ylim(0, 25)
    runtime_ax.set_ylabel("完整运行时间 (min)", fontproperties=font, color="#EA4335")
    runtime_ax.spines[["top", "right", "left"]].set_visible(False)
    runtime_ax.tick_params(axis="y", length=0, colors="#EA4335")
    for index, value in enumerate(runtime):
        runtime_ax.text(
            index, value + 0.9, f"{value:.1f}", ha="center", fontsize=9, color="#EA4335"
        )
    ax.set_title(
        "同一主机 · 同一 CQ-125 · 26B 是准确率与时延的工程平衡点",
        fontproperties=font,
        fontsize=15,
        fontweight="bold",
        color="#161D2C",
        loc="left",
        pad=14,
    )
    fig.tight_layout(pad=1.4)
    fig.savefig(output, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def build_farmland_optimization_map(
    background_shapefile: Path,
    changes_geojson: Path,
    summary: dict[str, Any],
    output: Path,
) -> None:
    background = gpd.read_file(background_shapefile, columns=[])
    changes = gpd.read_file(changes_geojson, columns=["CHG_FLAG"])
    if background.crs != changes.crs:
        changes = changes.to_crs(background.crs)
    farm_to_forest = changes[changes["CHG_FLAG"] == 1]
    forest_to_farm = changes[changes["CHG_FLAG"] == 2]
    spatial_output = summary["shapefile_output"]
    expected_swaps = int(summary["results"][0]["swaps_completed"])
    if len(farm_to_forest) != expected_swaps or len(forest_to_farm) != expected_swaps:
        raise ValueError("Change-map feature counts disagree with the planning summary")
    unchanged_count = int(spatial_output["n_input"]) - len(changes)

    fig, ax = plt.subplots(figsize=(9, 9), dpi=170)
    fig.patch.set_facecolor("#F7F9FC")
    ax.set_facecolor("#F7F9FC")
    background.plot(ax=ax, color="#D8DEE8", edgecolor="none", rasterized=True)
    farm_to_forest.plot(ax=ax, color="#EA4335", edgecolor="#B82C20", linewidth=0.12)
    forest_to_farm.plot(ax=ax, color="#34A853", edgecolor="#1F7A38", linewidth=0.12)
    ax.set_axis_off()
    ax.margins(0.015)
    ax.legend(
        handles=[
            Patch(facecolor="#D8DEE8", label=f"保持不变 {unchanged_count:,}"),
            Patch(facecolor="#EA4335", label=f"耕地 → 林地 {len(farm_to_forest):,}"),
            Patch(facecolor="#34A853", label=f"林地 → 耕地 {len(forest_to_farm):,}"),
        ],
        loc="lower left",
        prop=chart_font(),
        frameon=True,
        framealpha=0.96,
        facecolor="white",
        edgecolor="#DBE2EB",
        fontsize=10,
    )
    fig.tight_layout(pad=0.15)
    fig.savefig(output, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def build_nl2sql_crop(source: Path, output: Path) -> None:
    with Image.open(source) as image:
        if image.width / image.height > 1.4:
            rgb = image.convert("RGB")
            result_crop = rgb.crop(
                (
                    int(image.width * 0.047),
                    int(image.height * 0.255),
                    int(image.width * 0.214),
                    int(image.height * 0.785),
                )
            )
            map_crop = rgb.crop(
                (
                    int(image.width * 0.214),
                    int(image.height * 0.037),
                    int(image.width * 0.625),
                    image.height,
                )
            )
            canvas = Image.new("RGB", (1500, 1500), "#E8EAED")
            result_panel = ImageOps.contain(
                result_crop,
                (650, 1500),
                method=Image.Resampling.LANCZOS,
            )
            map_panel = ImageOps.contain(
                map_crop,
                (850, 1500),
                method=Image.Resampling.LANCZOS,
            )
            canvas.paste(result_panel, (0, int((1500 - result_panel.height) / 2)))
            canvas.paste(map_panel, (650, int((1500 - map_panel.height) / 2)))
            fitted = canvas
        else:
            crop = image.crop((70, 250, 1320, 1770)).convert("RGB")
            fitted = ImageOps.fit(crop, (1500, 1500), method=Image.Resampling.LANCZOS)
        fitted.save(output, quality=94)


def set_fill(shape, fill: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)


def set_line(shape, line: str = "line", width: float = 0.8) -> None:
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(width)


def apply_font(run, size: float, fill: str = "ink", bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color(fill)
    run.font.bold = bold


def add_text(
    slide,
    x,
    y,
    w,
    h,
    value: str,
    size: float = 16,
    fill: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = anchor
    for index, line in enumerate(value.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        run = paragraph.add_run()
        run.text = line
        apply_font(run, size, fill, bold)
    return box


def rect(slide, x, y, w, h, fill: str = "white", line: str = "line", radius: bool = False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    set_fill(shape, fill)
    set_line(shape, line)
    return shape


def pill(slide, x, y, w, h, label: str, fill: str, text_fill: str = "white", size: float = 11):
    shape = rect(slide, x, y, w, h, fill, fill, radius=True)
    add_text(slide, x, y, w, h, label, size, text_fill, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return shape


def add_image(slide, path: Path, x, y, w, h, line: str = "line"):
    rect(slide, x, y, w, h, "white", line, radius=True)
    with Image.open(path) as image:
        width, height = image.size
    ratio = min(w / width, h / height)
    placed_w = int(width * ratio)
    placed_h = int(height * ratio)
    placed_x = x + int((w - placed_w) / 2)
    placed_y = y + int((h - placed_h) / 2)
    return slide.shapes.add_picture(str(path), placed_x, placed_y, placed_w, placed_h)


def add_notes(slide, text: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    frame.text = text


def base_background(slide, dark: bool = False) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color("navy" if dark else "bg")
    x = 0.0
    for fill, width in (("blue", 3.333), ("red", 3.333), ("yellow", 3.333), ("green", 3.333)):
        bar = rect(slide, Inches(x), 0, Inches(width), Inches(0.065), fill, fill)
        bar.line.fill.background()
        x += width


def header(slide, section: str, title: str, subtitle: str, page: str) -> None:
    base_background(slide)
    pill(
        slide,
        Inches(0.56),
        Inches(0.27),
        Inches(1.42),
        Inches(0.32),
        section,
        "pale_blue",
        "blue",
        9.5,
    )
    add_text(slide, Inches(0.56), Inches(0.72), Inches(11.6), Inches(0.45), title, 23, "ink", True)
    add_text(slide, Inches(0.58), Inches(1.18), Inches(11.3), Inches(0.3), subtitle, 11.5, "muted")
    rect(slide, Inches(0.58), Inches(1.6), Inches(0.62), Inches(0.045), "green", "green")
    footer(slide, page)


def footer(slide, page: str, source: str | None = None) -> None:
    add_text(
        slide,
        Inches(0.56),
        Inches(7.08),
        Inches(5.4),
        Inches(0.18),
        "GIS Data Agent · Gemma 4 AI Agent Finals",
        8,
        "muted",
    )
    if source:
        add_text(
            slide,
            Inches(5.0),
            Inches(7.08),
            Inches(6.9),
            Inches(0.18),
            source,
            7.2,
            "muted",
            False,
            PP_ALIGN.RIGHT,
        )
    add_text(
        slide,
        Inches(12.22),
        Inches(7.06),
        Inches(0.55),
        Inches(0.2),
        page,
        8.5,
        "muted",
        False,
        PP_ALIGN.RIGHT,
    )


def metric_block(slide, x, y, w, value: str, label: str, note: str, accent: str):
    rect(slide, x, y, w, Inches(1.08), "white", "line", radius=True)
    rect(slide, x, y, Inches(0.055), Inches(1.08), accent, accent)
    add_text(
        slide,
        x + Inches(0.18),
        y + Inches(0.12),
        w - Inches(0.32),
        Inches(0.32),
        value,
        22,
        accent,
        True,
    )
    add_text(
        slide,
        x + Inches(0.18),
        y + Inches(0.51),
        w - Inches(0.32),
        Inches(0.2),
        label,
        10.5,
        "ink",
        True,
    )
    add_text(
        slide,
        x + Inches(0.18),
        y + Inches(0.77),
        w - Inches(0.32),
        Inches(0.18),
        note,
        8.2,
        "muted",
    )


def flow_arrow(slide, x1, y1, x2, y2, fill: str = "muted"):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color(fill)
    line.line.width = Pt(1.6)
    line.line.end_arrowhead = True
    return line


def slide_cover(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_background(slide, dark=True)
    rect(slide, Inches(0.0), Inches(0.065), Inches(7.35), Inches(7.435), "bg", "bg")
    add_text(
        slide,
        Inches(0.72),
        Inches(0.86),
        Inches(5.9),
        Inches(0.55),
        "GIS Data Agent",
        34,
        "ink",
        True,
    )
    add_text(
        slide,
        Inches(0.74),
        Inches(1.45),
        Inches(5.95),
        Inches(0.36),
        "Gemma 4 驱动的受控空间决策 Agent",
        17,
        "blue",
        True,
    )
    add_text(
        slide,
        Inches(0.74),
        Inches(2.07),
        Inches(5.9),
        Inches(1.05),
        "让自然语言问题落到\n可执行、可审计、可复用的 GIS 行动",
        23,
        "ink",
        True,
    )
    stages = [("问", "blue"), ("算", "green"), ("验", "yellow"), ("记", "red")]
    for index, (label, accent) in enumerate(stages):
        x = Inches(0.78 + index * 1.35)
        pill(slide, x, Inches(3.46), Inches(0.76), Inches(0.55), label, accent, "white", 16)
        if index < len(stages) - 1:
            add_text(
                slide,
                x + Inches(0.85),
                Inches(3.56),
                Inches(0.26),
                Inches(0.22),
                "→",
                14,
                "muted",
                True,
            )
    value = evidence["result"]
    mini = [
        ("90.4%", "CQ-125", "blue"),
        ("30/30", "Gemma 4 + ADK", "green"),
        ("硬约束校验", "真实 MPC 规划", "red"),
    ]
    for index, (number, label, accent) in enumerate(mini):
        x = Inches(0.74 + index * 2.02)
        rect(slide, x, Inches(4.65), Inches(1.72), Inches(0.94), "white", "line", radius=True)
        add_text(
            slide,
            x + Inches(0.15),
            Inches(4.8),
            Inches(1.42),
            Inches(0.26),
            number,
            15,
            accent,
            True,
            PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x + Inches(0.15),
            Inches(5.18),
            Inches(1.42),
            Inches(0.18),
            label,
            8.5,
            "muted",
            False,
            PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        Inches(0.76),
        Inches(6.45),
        Inches(5.7),
        Inches(0.2),
        "周宁 · Gemma 4 AI Agent 赛道 A 决赛",
        9.5,
        "muted",
    )

    add_image(slide, assets["map"], Inches(7.66), Inches(0.74), Inches(5.12), Inches(5.72), "navy")
    pill(
        slide,
        Inches(8.0),
        Inches(5.98),
        Inches(2.15),
        Inches(0.42),
        "101,657 条空间记录",
        "white",
        "navy",
        9.5,
    )
    pill(
        slide,
        Inches(10.36),
        Inches(5.98),
        Inches(2.03),
        Inches(0.42),
        f"坡度 {value['slope_change_pct']:.3f}%",
        "green",
        "white",
        9.5,
    )
    add_text(
        slide,
        Inches(7.82),
        Inches(6.72),
        Inches(4.75),
        Inches(0.22),
        "县域耕地空间优化引擎 · 真实 MPC 规划产物",
        9,
        "white",
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(12.18),
        Inches(7.05),
        Inches(0.6),
        Inches(0.2),
        "01",
        8.5,
        "white",
        False,
        PP_ALIGN.RIGHT,
    )
    add_notes(
        slide,
        (
            "0:00-0:30。开场：GIS Data Agent 不是会聊天的地图，而是把自然语言问题落到"
            "可执行、可审计、可复用 GIS 行动的受控 Agent。今天只展示两件事：高频空间"
            "问数，以及高价值县域规划。右侧不是示意图，是本次决赛版本真实 MPC 规划产物。"
        ),
    )
    return slide


def slide_problem_solution(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del evidence, assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(
        slide,
        "真实影响 30%",
        "空间数据不是问答题，而是带约束的执行题",
        "同一个 Agent 覆盖高频问数与低频高价值规划，但每一步都必须留下证据。",
        "02",
    )

    rect(
        slide, Inches(0.65), Inches(1.94), Inches(5.78), Inches(2.52), "white", "line", radius=True
    )
    pill(
        slide,
        Inches(0.92),
        Inches(2.2),
        Inches(1.24),
        Inches(0.34),
        "高频刚需",
        "blue",
        "white",
        9.5,
    )
    add_text(
        slide,
        Inches(0.92),
        Inches(2.73),
        Inches(4.95),
        Inches(0.35),
        "NL2Semantic2GeoSQL",
        18,
        "ink",
        True,
    )
    add_text(
        slide,
        Inches(0.92),
        Inches(3.2),
        Inches(5.0),
        Inches(0.64),
        (
            "空间距离、面积、SRID、geometry/geography、空间 join 去重，都会让普通 "
            "NL2SQL 的“看似正确”在执行时失败。"
        ),
        12.2,
        "muted",
    )
    add_text(
        slide,
        Inches(0.92),
        Inches(4.02),
        Inches(5.0),
        Inches(0.2),
        "用户：GIS 分析师、空间数据库团队",
        10.5,
        "blue",
        True,
    )

    rect(slide, Inches(6.9), Inches(1.94), Inches(5.78), Inches(2.52), "white", "line", radius=True)
    pill(
        slide,
        Inches(7.17),
        Inches(2.2),
        Inches(1.52),
        Inches(0.34),
        "高价值决策",
        "green",
        "white",
        9.5,
    )
    add_text(
        slide,
        Inches(7.17),
        Inches(2.73),
        Inches(4.95),
        Inches(0.35),
        "县域耕地空间布局优化",
        18,
        "ink",
        True,
    )
    add_text(
        slide,
        Inches(7.17),
        Inches(3.2),
        Inches(5.0),
        Inches(0.64),
        "结果必须同时满足耕地面积不减少、坡度下降、连片度提升，并生成可交付的空间成果，而不是文字建议。",
        12.2,
        "muted",
    )
    add_text(
        slide,
        Inches(7.17),
        Inches(4.02),
        Inches(5.0),
        Inches(0.2),
        "用户：自然资源、国土规划、县域整治人员",
        10.5,
        "green",
        True,
    )

    flow_y = Inches(5.18)
    steps = [
        ("目标", "Gemma 4", "blue"),
        ("运行时", "Google ADK", "green"),
        ("行动", "PostGIS / MPC 引擎", "yellow"),
        ("校验", "硬约束校验", "red"),
        ("复用", "已验证经验库", "cyan"),
    ]
    for index, (kind, label, accent) in enumerate(steps):
        x = Inches(0.62 + index * 2.54)
        rect(slide, x, flow_y, Inches(2.03), Inches(0.86), "white", "line", radius=True)
        add_text(
            slide,
            x + Inches(0.14),
            flow_y + Inches(0.12),
            Inches(1.75),
            Inches(0.16),
            kind,
            8.5,
            accent,
            True,
            PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x + Inches(0.14),
            flow_y + Inches(0.42),
            Inches(1.75),
            Inches(0.2),
            label,
            11.2,
            "ink",
            True,
            PP_ALIGN.CENTER,
        )
        if index < len(steps) - 1:
            flow_arrow(
                slide,
                x + Inches(2.08),
                flow_y + Inches(0.43),
                x + Inches(2.41),
                flow_y + Inches(0.43),
            )
    add_text(
        slide,
        Inches(1.0),
        Inches(6.38),
        Inches(11.3),
        Inches(0.32),
        "核心边界：模型负责选择和解释，确定性工具负责计算，治理代码决定结果能否保存和复用。",
        13,
        "navy",
        True,
        PP_ALIGN.CENTER,
    )
    add_notes(
        slide,
        (
            "0:30-1:05。空间问数是高频刚需，难点不是生成 SQL，而是空间语义与真实执行。"
            "县域规划频率更低，但每次价值更高，必须满足业务硬约束并交付空间成果。架构上 "
            "Gemma 4 做决策，ADK 承载原生函数调用，PostGIS 与县域耕地空间优化引擎"
            "做确定性计算，硬约束校验决定结果能否进入已验证经验库。"
        ),
    )
    return slide


def slide_nl2sql(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    selected = next(row for row in evidence["benchmark"] if row["model"] == "Gemma4:26b")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(
        slide,
        "Demo 1 · 0:50",
        "NL2Semantic2GeoSQL：不是生成 SQL，而是让空间语义真正执行",
        "同一 CQ-125 基准选模型；现场锁定一个问题、一段 GeoSQL、一个可核查地图结果。",
        "03",
    )
    add_image(
        slide, assets["nl2sql"], Inches(0.62), Inches(1.86), Inches(4.45), Inches(4.92), "line"
    )
    add_image(
        slide, assets["benchmark"], Inches(5.42), Inches(1.84), Inches(7.25), Inches(2.78), "line"
    )

    rect(
        slide, Inches(5.42), Inches(4.91), Inches(7.25), Inches(1.85), "white", "line", radius=True
    )
    pill(
        slide,
        Inches(5.7),
        Inches(5.16),
        Inches(1.2),
        Inches(0.34),
        "锁定问题",
        "blue",
        "white",
        9.5,
    )
    add_text(
        slide,
        Inches(7.08),
        Inches(5.08),
        Inches(5.1),
        Inches(0.38),
        "最长桥梁 100 米范围内有多少高德 POI？",
        14,
        "ink",
        True,
    )
    add_text(
        slide,
        Inches(5.7),
        Inches(5.56),
        Inches(6.45),
        Inches(0.24),
        "语义 / Schema Grounding → Gemma 4 生成 GeoSQL → 只读执行 → 地图核验",
        9.6,
        "blue",
        True,
    )
    add_text(
        slide,
        Inches(5.7),
        Inches(5.95),
        Inches(4.6),
        Inches(0.27),
        "ST_DWithin(...::geography, 100)",
        12.5,
        "navy",
        True,
    )
    pill(
        slide,
        Inches(10.62),
        Inches(5.84),
        Inches(1.55),
        Inches(0.52),
        "35 + 地图",
        "green",
        "white",
        15,
    )
    add_text(
        slide,
        Inches(5.7),
        Inches(6.4),
        Inches(6.45),
        Inches(0.2),
        (
            f"Gemma 4 26B：{selected['full_ex_correct']}/{selected['full_ex_total']} = "
            f"{float(selected['full_ex_pct']):.1f}% · 只比 31B 少 1 题，约快 37.6%"
        ),
        9.5,
        "muted",
    )
    footer(slide, "03", "证据：CQ-125 模型选型 + 锁定 PostGIS 演示快照")
    add_notes(
        slide,
        (
            "1:05-1:55。这里先播放约 35 秒锁定视频，或直接展示结果。问题是最长桥梁 100 米"
            "范围内有多少高德 POI。系统先做语义与 schema grounding，再由 Gemma 4 生成 GeoSQL，"
            "关键是 geography 米制距离和只读执行；结果为 35，并自动加载同一 PostGIS 查询快照的"
            "桥梁、百米范围和 35 个 POI。模型不是凭印象选的：同一 CQ-125 上 "
            "26B 为 113/125，只比 31B 少一题，但完整运行约快 37.6%。"
        ),
    )
    return slide


def slide_farmland_planning(
    prs: Presentation,
    evidence: dict[str, Any],
    assets: dict[str, Path],
):
    result = evidence["result"]
    summary = evidence["summary"]
    spatial_output = summary["shapefile_output"]
    config = summary["config"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(
        slide,
        "Demo 2 · 1:30",
        "Gemma 4 + 地理空间世界模型领域原型",
        (
            "县域耕地空间优化引擎预测行动后果并搜索方案；Gemma 4 + ADK 负责目标理解、"
            "工具决策与闭环控制。"
        ),
        "04",
    )
    add_image(slide, assets["map"], Inches(0.62), Inches(1.84), Inches(5.18), Inches(4.95), "line")

    trace = [
        ("1", "状态检查", "观察版本", "blue"),
        ("2", "资源检查", "判断复用", "green"),
        ("3", "经验召回", "注入上下文", "cyan"),
        ("4", "MPC 规划", "执行空间行动", "yellow"),
        ("5", "硬约束审计", "独立验收", "red"),
        ("6", "经验写入 / 人工", "按结果分支", "green"),
    ]
    pill(
        slide,
        Inches(6.18),
        Inches(1.84),
        Inches(3.86),
        Inches(0.34),
        "LLM 控制面 · Gemma 4 + ADK · 6 次函数调用",
        "pale_blue",
        "blue",
        8.8,
    )
    for index, (number, tool, meaning, accent) in enumerate(trace):
        y = Inches(2.35 + index * 0.43)
        pill(slide, Inches(6.18), y, Inches(0.4), Inches(0.3), number, accent, "white", 9)
        add_text(
            slide,
            Inches(6.75),
            y - Inches(0.01),
            Inches(1.82),
            Inches(0.22),
            tool,
            10.3,
            "ink",
            True,
        )
        add_text(
            slide,
            Inches(8.7),
            y - Inches(0.01),
            Inches(1.34),
            Inches(0.22),
            meaning,
            8.7,
            "muted",
        )
        if index < len(trace) - 1:
            flow_arrow(
                slide,
                Inches(6.38),
                y + Inches(0.31),
                Inches(6.38),
                y + Inches(0.4),
                accent,
            )

    rect(
        slide, Inches(10.38), Inches(1.95), Inches(2.27), Inches(2.95), "navy", "navy", radius=True
    )
    add_text(
        slide,
        Inches(10.62),
        Inches(2.2),
        Inches(1.78),
        Inches(0.24),
        "真实 MPC 规划",
        12,
        "white",
        True,
        PP_ALIGN.CENTER,
    )
    metrics = [
        (f"{spatial_output['n_input']:,}", "输入空间记录"),
        (f"{spatial_output['n_in_env']:,}", "环境图斑"),
        (f"{config['n_blocks']:,}", "空间块"),
        (f"{result['swaps_completed']:,}", "双向置换"),
    ]
    for index, (value, label) in enumerate(metrics):
        y = Inches(2.72 + index * 0.48)
        add_text(
            slide,
            Inches(10.58),
            y,
            Inches(0.84),
            Inches(0.22),
            value,
            12.5,
            "white",
            True,
            PP_ALIGN.RIGHT,
        )
        add_text(
            slide, Inches(11.55), y + Inches(0.02), Inches(0.78), Inches(0.18), label, 7.9, "white"
        )

    rect(
        slide,
        Inches(6.15),
        Inches(5.02),
        Inches(6.5),
        Inches(0.37),
        "pale_yellow",
        "yellow",
        radius=True,
    )
    add_text(
        slide,
        Inches(6.3),
        Inches(5.1),
        Inches(6.2),
        Inches(0.18),
        "GWM 领域原型 · S 空间状态 + A 候选置换 → Tθ 转移预测 → MPC 搜索 → 校验",
        9.2,
        "navy",
        True,
        PP_ALIGN.CENTER,
    )

    changes = [
        (f"+{result['cultivated_area_change_ha']:.3f} ha", "耕地面积", "green"),
        (f"{result['slope_change_pct']:.3f}%", "平均坡度", "blue"),
        (f"+{result['cont_change']:.4f}", "连片度", "cyan"),
        (f"+{result['baimu_area_change_ha']:.2f} ha", "百亩方", "red"),
    ]
    for index, (value, label, accent) in enumerate(changes):
        x = Inches(6.15 + (index % 2) * 3.3)
        y = Inches(5.55 + (index // 2) * 0.66)
        rect(slide, x, y, Inches(3.02), Inches(0.58), "white", "line", radius=True)
        add_text(
            slide,
            x + Inches(0.16),
            y + Inches(0.11),
            Inches(1.6),
            Inches(0.23),
            value,
            12,
            accent,
            True,
        )
        add_text(
            slide,
            x + Inches(1.84),
            y + Inches(0.13),
            Inches(0.92),
            Inches(0.2),
            label,
            9,
            "muted",
            False,
            PP_ALIGN.RIGHT,
        )
    footer(
        slide,
        "04",
        "当前边界：领域化、基于模型的空间规划原型，不等同于完整通用 GWM",
    )
    add_notes(
        slide,
        (
            "1:55-3:25。现场运行成功链或播放约 70 秒同版本视频。Gemma 4 通过 Google ADK"
            "原生函数调用检查版本与资源、召回经验，并根据观察选择快速 MPC 路径；硬约束审计"
            "通过才写入已验证经验库，失败只允许重规划一次，再失败转人工。空间算法使用学习型"
            "状态转移模型集成预测候选行动，再由 MPC 搜索，确定性 Python 负责数值计算。Paper9"
            "不是完整 GWM，而是它的领域原型：图斑与空间块构成状态，地类置换构成行动，学习型"
            "转移模型预测行动后果，MPC 消费这些预测并搜索方案。县域耕地空间优化引擎的"
            "内部研发代号是 Paper9，但主路演统一使用业务名称。"
            f"本次真实运行处理 {spatial_output['n_input']:,} 条输入、"
            f"{spatial_output['n_in_env']:,} 个环境图斑、{config['n_blocks']:,} 个空间块，"
            f"完成 {result['swaps_completed']:,} 对置换；耕地面积不减少，坡度下降，"
            "连片度和百亩方面积提升。A/B/C/D 是算法阶段，不冒充 Gemma 4 动态规划。"
        ),
    )
    return slide


def slide_evidence(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del assets
    passed = evidence["reliability_passed"]
    total = evidence["reliability_total"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(
        slide,
        "证据 · 1:00",
        "这不是一次成功 Demo，而是一套分层证据",
        "模型编排、算法运行、工程质量和历史环境记录分别证明不同问题。",
        "05",
    )
    blocks = [
        (f"{passed}/{total}", "Gemma 4 + ADK", "三场景各 10 次", "blue"),
        ("236", "镜像运行时测试", "另有 6 交付 + 52 兼容", "green"),
        ("3/3", "容器 healthy", "app / PostGIS / Redis", "yellow"),
        ("硬约束校验", "真实 MPC 规划", "1 次规划 + 空间产物", "red"),
    ]
    for index, block in enumerate(blocks):
        metric_block(slide, Inches(0.66 + index * 3.12), Inches(1.92), Inches(2.72), *block)

    rect(slide, Inches(0.66), Inches(3.36), Inches(7.3), Inches(2.7), "white", "line", radius=True)
    add_text(
        slide,
        Inches(0.96),
        Inches(3.62),
        Inches(6.65),
        Inches(0.28),
        "可靠性基线：成功、恢复、停止都可重复",
        15,
        "ink",
        True,
    )
    rows = [
        ("首次通过", "6 次函数调用", "10/10", "7.97 s"),
        ("版本不兼容", "2 次调用 → 停止", "10/10", "3.24 s"),
        ("一次重规划", "8 次函数调用", "10/10", "13.57 s"),
    ]
    for index, row in enumerate(rows):
        y = Inches(4.16 + index * 0.47)
        if index % 2 == 0:
            rect(
                slide,
                Inches(0.92),
                y - Inches(0.04),
                Inches(6.76),
                Inches(0.38),
                "pale_blue",
                "pale_blue",
            )
        for col, (value, width) in enumerate(zip(row, (1.5, 1.85, 1.2, 1.2), strict=True)):
            x = Inches(1.02 + sum((1.5, 1.85, 1.2, 1.2)[:col]))
            add_text(slide, x, y, Inches(width), Inches(0.2), value, 10.2, "ink", col == 0)
    add_text(
        slide,
        Inches(0.96),
        Inches(5.7),
        Inches(6.65),
        Inches(0.2),
        "30/30 Wilson 95% CI：88.65%–100% · 工具响应为确定性替身",
        8.8,
        "muted",
    )

    rect(
        slide,
        Inches(8.3),
        Inches(3.36),
        Inches(4.38),
        Inches(2.7),
        "pale_yellow",
        "yellow",
        radius=True,
    )
    add_text(
        slide,
        Inches(8.62),
        Inches(3.63),
        Inches(3.72),
        Inches(0.3),
        "真实环境可行性：证据分层",
        14,
        "ink",
        True,
    )
    add_text(
        slide,
        Inches(8.62),
        Inches(4.18),
        Inches(3.72),
        Inches(0.96),
        (
            "历史版本：内部交接材料记录在目标内网使用真实权威数据完成全流程。\n"
            "当前版本：0.3.3 / 2.2.3 本机 MPC 完成；\n硬约束校验通过。"
        ),
        11.2,
        "ink",
    )
    add_text(
        slide,
        Inches(8.62),
        Inches(5.34),
        Inches(3.72),
        Inches(0.4),
        "历史记录 ≠ 当前部内复测 ≠ GIS Data Agent 生产部署",
        9.5,
        "muted",
        True,
    )
    footer(slide, "05", "证据分类相互独立；详见 claim_register.md")
    add_notes(
        slide,
        (
            "3:25-4:20。这不是一次跑通就算完成。Gemma 4 与 ADK 在首次成功、版本阻断和"
            "一次恢复三个分支各跑十次，30/30 通过；Wilson 区间说明我们没有把它夸大成"
            "普遍 100%。工程侧分为 236 个镜像运行时测试、6 个宿主交付契约、52 个兼容"
            "测试和三个健康容器；这些范围不同，不能混成一个总数。真实"
            "环境证据也严格分层：历史内网记录证明底层耕地优化引擎的工程可行性，当前 "
            "v2.2.3 本机运行证明适配与 MPC 规划，不把两者混成生产验收。"
        ),
    )
    return slide


def slide_close(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del evidence, assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    base_background(slide, dark=True)
    add_text(
        slide,
        Inches(0.72),
        Inches(0.78),
        Inches(11.5),
        Inches(0.56),
        "下一代空间智能体：LLM + GWM",
        30,
        "white",
        True,
    )
    add_text(
        slide,
        Inches(0.74),
        Inches(1.46),
        Inches(11.2),
        Inches(0.38),
        "语言智能负责理解与协同，地理空间世界模型负责推演行动后果。",
        16.5,
        "white",
    )
    rect(slide, Inches(0.74), Inches(2.25), Inches(11.9), Inches(0.055), "green", "green")

    columns = [
        ("Gemma 4 + ADK", "理解目标\n选择工具\n人机协同", "blue"),
        ("GWM", "表示空间状态\n预测行动后果\n反事实推演", "green"),
        ("规划器", "MPC / 搜索\n方案比较\n行动选择", "yellow"),
        ("审核控制", "规则校验\n证据边界\n保存 / 人工", "red"),
    ]
    for index, (title, body, accent) in enumerate(columns):
        x = Inches(0.72 + index * 3.13)
        rect(slide, x, Inches(2.78), Inches(2.72), Inches(2.2), "white", "white", radius=True)
        pill(
            slide,
            x + Inches(0.25),
            Inches(3.04),
            Inches(1.55),
            Inches(0.35),
            title,
            accent,
            "white",
            9.5,
        )
        add_text(
            slide,
            x + Inches(0.28),
            Inches(3.68),
            Inches(2.15),
            Inches(0.92),
            body,
            12.2,
            "navy",
            True,
            PP_ALIGN.LEFT,
        )
        if index < len(columns) - 1:
            add_text(
                slide,
                x + Inches(2.78),
                Inches(3.72),
                Inches(0.28),
                Inches(0.3),
                "→",
                16,
                "white",
                True,
                PP_ALIGN.CENTER,
            )

    rect(
        slide,
        Inches(0.88),
        Inches(5.43),
        Inches(11.58),
        Inches(0.76),
        "white",
        "white",
        radius=True,
    )
    add_text(
        slide,
        Inches(1.12),
        Inches(5.58),
        Inches(2.0),
        Inches(0.22),
        "今天已验证",
        12,
        "green",
        True,
    )
    add_text(
        slide,
        Inches(2.55),
        Inches(5.56),
        Inches(9.55),
        Inches(0.28),
        "空间问数 + GWM 领域原型 + 硬约束校验 + 已验证经验",
        12,
        "navy",
        True,
    )
    add_text(
        slide,
        Inches(1.12),
        Inches(6.45),
        Inches(11.1),
        Inches(0.26),
        "路线：多尺度状态 · 多行动动力学 · 长期/不确定性推演 · 因果校准 · 证据校验",
        12.5,
        "white",
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(12.18),
        Inches(7.05),
        Inches(0.6),
        Inches(0.2),
        "06",
        8.5,
        "white",
        False,
        PP_ALIGN.RIGHT,
    )
    add_notes(
        slide,
        (
            "4:20-4:40。总结：我们的目标不是不断给 LLM 增加 GIS 工具，而是走向 LLM + GWM。"
            "Gemma 4 与 ADK 负责理解目标、选择工具和人机协同；GWM 负责表示空间状态、预测"
            "行动后果和反事实推演；MPC 消费预测做方案选择，审核机制决定能否保存。今天已经"
            "验证空间问数、GWM 领域原型和受控闭环，下一步扩展多尺度、多行动、长期与"
            "不确定性推演。谢谢。"
        ),
    )
    return slide


def appendix_header(slide, title: str, subtitle: str, page: str):
    header(slide, "Q&A 附录", title, subtitle, page)


def appendix_architecture(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del evidence, assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    appendix_header(
        slide,
        "Gemma 4 真的在规划，还是 Python 固定流程？",
        "直接结论：Gemma 4 规划工具调用与恢复分支；确定性代码负责空间数值计算。",
        "A1",
    )
    rect(slide, Inches(0.72), Inches(1.88), Inches(5.9), Inches(3.55), "white", "line", radius=True)
    pill(
        slide,
        Inches(0.98),
        Inches(2.12),
        Inches(2.25),
        Inches(0.38),
        "Agent 规划 · Gemma 4 + ADK",
        "blue",
        "white",
        9.3,
    )
    add_text(
        slide,
        Inches(4.95),
        Inches(2.17),
        Inches(1.35),
        Inches(0.2),
        "模型动态选择",
        9.5,
        "blue",
        True,
        PP_ALIGN.RIGHT,
    )
    agent_steps = ["理解目标", "版本资源\n检查", "召回经验", "选择工具", "观察后分支"]
    for index, label in enumerate(agent_steps):
        x = Inches(0.98 + index * 1.08)
        rect(slide, x, Inches(3.0), Inches(0.88), Inches(0.88), "pale_blue", "blue", radius=True)
        add_text(
            slide,
            x + Inches(0.08),
            Inches(3.19),
            Inches(0.72),
            Inches(0.44),
            label,
            8.6,
            "blue",
            True,
            PP_ALIGN.CENTER,
            MSO_ANCHOR.MIDDLE,
        )
        if index < len(agent_steps) - 1:
            flow_arrow(
                slide,
                x + Inches(0.89),
                Inches(3.44),
                x + Inches(1.03),
                Inches(3.44),
                "blue",
            )
    add_text(
        slide,
        Inches(1.0),
        Inches(4.33),
        Inches(5.3),
        Inches(0.54),
        "可停止 · 可保存 · 可重规划一次 · 可转人工",
        11.5,
        "ink",
        True,
        PP_ALIGN.CENTER,
    )

    rect(
        slide,
        Inches(6.82),
        Inches(1.88),
        Inches(5.78),
        Inches(3.55),
        "white",
        "line",
        radius=True,
    )
    pill(
        slide,
        Inches(7.08),
        Inches(2.12),
        Inches(2.7),
        Inches(0.38),
        "算法规划 · 县域耕地优化引擎",
        "green",
        "white",
        9.3,
    )
    add_text(
        slide,
        Inches(11.08),
        Inches(2.17),
        Inches(1.22),
        Inches(0.2),
        "确定性执行",
        9.5,
        "green",
        True,
        PP_ALIGN.RIGHT,
    )
    algorithm_steps = [
        ("S", "空间状态"),
        ("Tθ", "预测后果"),
        ("MPC", "搜索方案"),
        ("校验", "验收成果"),
    ]
    for index, (symbol, label) in enumerate(algorithm_steps):
        x = Inches(7.15 + index * 1.33)
        rect(slide, x, Inches(2.93), Inches(1.05), Inches(1.02), "pale_green", "green", radius=True)
        add_text(
            slide,
            x,
            Inches(3.07),
            Inches(1.05),
            Inches(0.24),
            symbol,
            13,
            "green",
            True,
            PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            x,
            Inches(3.48),
            Inches(1.05),
            Inches(0.22),
            label,
            8.8,
            "ink",
            True,
            PP_ALIGN.CENTER,
        )
        if index < len(algorithm_steps) - 1:
            flow_arrow(
                slide,
                x + Inches(1.06),
                Inches(3.44),
                x + Inches(1.27),
                Inches(3.44),
                "green",
            )
    add_text(
        slide,
        Inches(7.12),
        Inches(4.33),
        Inches(5.1),
        Inches(0.54),
        "A/B/C/D 是算法阶段，不冒充 LLM 动态规划",
        11.5,
        "ink",
        True,
        PP_ALIGN.CENTER,
    )

    branches = [
        ("2 次调用", "版本不兼容 → 停止", "red"),
        ("6 次调用", "首次审计通过 → 保存", "green"),
        ("8 次调用", "审计失败 → 一次恢复", "yellow"),
    ]
    for index, (count, result, accent) in enumerate(branches):
        x = Inches(0.82 + index * 4.16)
        rect(slide, x, Inches(5.78), Inches(3.72), Inches(0.76), "white", "line", radius=True)
        pill(
            slide,
            x + Inches(0.18),
            Inches(5.99),
            Inches(0.9),
            Inches(0.34),
            count,
            accent,
            "white",
            8.5,
        )
        add_text(
            slide,
            x + Inches(1.25),
            Inches(6.02),
            Inches(2.22),
            Inches(0.24),
            result,
            9.3,
            "ink",
            True,
        )
    footer(slide, "A1", "边界：Gemma 4 决策工具与分支；GIS 数值由确定性工具计算")
    add_notes(
        slide,
        (
            "结论先说：两层规划。Gemma 4 通过 ADK 根据版本、资源、经验和审计反馈选择"
            "下一工具与停止条件；县域耕地优化引擎则用状态转移模型和 MPC 做空间数值规划。"
            "三个 2/6/8 次调用分支已分别重复验证，但 A/B/C/D 仍是确定性算法阶段。"
        ),
    )
    return slide


def appendix_reliability(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del evidence, assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    appendix_header(
        slide,
        "县域耕地优化引擎是不是地理空间世界模型？",
        "直接结论：它是领域化 GWM 原型，不是完整通用地理空间世界模型。",
        "A2",
    )
    stages = [
        ("空间状态 S", "图斑 · 空间块\n地类 · 坡度", "blue"),
        ("转移模型 Tθ", "学习型模型集成\n预测行动后果", "green"),
        ("规划器 MPC", "消费预测\n比较候选方案", "yellow"),
        ("治理外环", "硬约束审计\n保存 / 转人工", "red"),
    ]
    for index, (title, body, accent) in enumerate(stages):
        x = Inches(0.72 + index * 3.13)
        rect(slide, x, Inches(2.02), Inches(2.72), Inches(1.78), "white", "line", radius=True)
        pill(
            slide,
            x + Inches(0.2),
            Inches(2.22),
            Inches(1.48),
            Inches(0.36),
            title,
            accent,
            "white",
            9.2,
        )
        add_text(
            slide,
            x + Inches(0.2),
            Inches(2.88),
            Inches(2.28),
            Inches(0.58),
            body,
            10.7,
            "ink",
            True,
            PP_ALIGN.CENTER,
        )
        if index < len(stages) - 1:
            flow_arrow(
                slide,
                x + Inches(2.73),
                Inches(2.91),
                x + Inches(3.04),
                Inches(2.91),
                accent,
            )

    rect(
        slide,
        Inches(0.82),
        Inches(4.28),
        Inches(5.72),
        Inches(1.7),
        "pale_green",
        "green",
        radius=True,
    )
    add_text(
        slide,
        Inches(1.1),
        Inches(4.56),
        Inches(5.12),
        Inches(0.28),
        "今天已经具备",
        14,
        "green",
        True,
    )
    add_text(
        slide,
        Inches(1.1),
        Inches(5.05),
        Inches(5.12),
        Inches(0.56),
        "状态表示 · 学习型转移预测 · MPC 搜索\n硬约束校验 · 已验证经验闭环",
        10.8,
        "ink",
        True,
    )
    rect(
        slide,
        Inches(6.8),
        Inches(4.28),
        Inches(5.72),
        Inches(1.7),
        "pale_red",
        "red",
        radius=True,
    )
    add_text(
        slide,
        Inches(7.08),
        Inches(4.56),
        Inches(5.12),
        Inches(0.28),
        "距离完整 GWM 还缺",
        14,
        "red",
        True,
    )
    add_text(
        slide,
        Inches(7.08),
        Inches(5.05),
        Inches(5.12),
        Inches(0.56),
        "多尺度 / 多行动 · 长期推演与不确定性\n因果校准 · 跨区域验证",
        10.8,
        "ink",
        True,
    )
    footer(slide, "A2", "内部研发代号 Paper9；MPC 是世界模型的消费者，不是世界模型本身")
    add_notes(
        slide,
        (
            "严格地说，它不是完整 GWM，而是领域化原型。图斑和空间块构成状态，学习型"
            "状态转移模型集成预测候选行动的后果，这一部分是 GWM 动力学内核的雏形；MPC"
            "消费预测并搜索方案，硬约束校验负责验收。下一步还要补齐多尺度、多行动、长期"
            "不确定性推演、因果校准和跨区域验证。"
        ),
    )
    return slide


def appendix_claims(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    appendix_header(
        slide,
        "结果为什么可信，失败时怎么办？",
        "直接结论：模型编排、真实算法和约束校验分别验证；任何一层失败都不会冒充成功。",
        "A3",
    )
    result = evidence["result"]
    cards = [
        (
            "模型编排证据",
            "30/30",
            "首次成功 / 版本阻断 / 一次恢复\nWilson 95% CI：88.65%–100%",
            "blue",
        ),
        (
            "真实算法证据",
            f"{result['swaps_completed']} 对置换",
            "101,657 条输入 · 真实 MPC\n面积不减 · 坡度下降 · 连片度提升",
            "green",
        ),
        (
            "治理证据",
            "硬约束校验通过",
            "产物完整才保存经验\n失败 → 重规划一次 → 转人工",
            "red",
        ),
    ]
    for index, (title, value, body, accent) in enumerate(cards):
        x = Inches(0.72 + index * 4.18)
        rect(slide, x, Inches(2.0), Inches(3.72), Inches(3.15), "white", "line", radius=True)
        pill(
            slide,
            x + Inches(0.25),
            Inches(2.26),
            Inches(1.7),
            Inches(0.38),
            title,
            accent,
            "white",
            9.3,
        )
        add_text(
            slide,
            x + Inches(0.25),
            Inches(3.0),
            Inches(3.2),
            Inches(0.42),
            value,
            18,
            accent,
            True,
        )
        add_text(
            slide,
            x + Inches(0.25),
            Inches(3.82),
            Inches(3.2),
            Inches(0.8),
            body,
            10.5,
            "ink",
        )
    rect(
        slide,
        Inches(0.95),
        Inches(5.72),
        Inches(11.42),
        Inches(0.74),
        "pale_yellow",
        "yellow",
        radius=True,
    )
    add_text(
        slide,
        Inches(1.22),
        Inches(5.96),
        Inches(10.88),
        Inches(0.26),
        "边界：30/30 验证的是 Gemma 4 + ADK 编排；真实 MPC 运行另行验证，二者不混算。",
        11.2,
        "ink",
        True,
        PP_ALIGN.CENTER,
    )
    footer(slide, "A3", "任何版本不兼容、二次审计失败或空间产物缺失均安全停止")
    add_notes(
        slide,
        (
            "可信性不是一个数字。三十次真实 Gemma 4 加 ADK 运行验证工具编排；本机真实 MPC"
            "运行验证算法和空间产物；硬约束校验决定能否保存。三十次使用确定性工具替身，"
            "不等于三十次算法重跑。任何版本、审计或产物失败都会停止，最多恢复一次后转人工。"
        ),
    )
    return slide


def appendix_google(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    appendix_header(
        slide,
        "真实落地证据是什么，哪些还没有证明？",
        "直接结论：需求、规模、历史环境与当前运行均有证据；生产成效仍需继续量化。",
        "A4",
    )
    spatial_output = evidence["summary"]["shapefile_output"]
    items = [
        (
            "真实需求",
            "GIS 分析师的高频空间问数；自然资源与县域整治人员的高价值布局优化",
            "受众清晰",
            "blue",
        ),
        (
            "当前规模",
            f"{spatial_output['n_input']:,} 条输入空间记录；CQ-125 可复核空间 SQL 基准",
            "可量化",
            "green",
        ),
        (
            "历史环境",
            "底层引擎旧版本的内部交接材料记录：曾在目标内网以真实权威数据完成全流程",
            "工程可行性",
            "yellow",
        ),
        (
            "当前运行",
            "0.3.3 / 2.2.3 本机真实 MPC、地图、图文 PDF、硬约束校验和经验保存",
            "当前可运行",
            "red",
        ),
    ]
    for index, (kind, fact, meaning, accent) in enumerate(items):
        y = Inches(1.9 + index * 1.02)
        pill(slide, Inches(0.75), y, Inches(1.22), Inches(0.38), kind, accent, "white", 9.3)
        rect(
            slide,
            Inches(2.2),
            y - Inches(0.08),
            Inches(10.25),
            Inches(0.82),
            "white",
            "line",
            radius=True,
        )
        add_text(
            slide,
            Inches(2.47),
            y + Inches(0.05),
            Inches(7.92),
            Inches(0.35),
            fact,
            10.5,
            "ink",
            True,
        )
        add_text(
            slide,
            Inches(10.58),
            y + Inches(0.09),
            Inches(1.5),
            Inches(0.22),
            meaning,
            9.3,
            accent,
            True,
            PP_ALIGN.RIGHT,
        )
    rect(
        slide,
        Inches(0.95),
        Inches(6.02),
        Inches(11.42),
        Inches(0.58),
        "pale_red",
        "red",
        radius=True,
    )
    add_text(
        slide,
        Inches(1.22),
        Inches(6.18),
        Inches(10.88),
        Inches(0.24),
        "尚未证明：v2.2.3 部内复测或生产验收、全国泛化、公开用户数、节省工时与人工修订率。",
        10.5,
        "red",
        True,
        PP_ALIGN.CENTER,
    )
    footer(slide, "A4", "历史内网记录 ≠ 当前版本部内复测 ≠ GIS Data Agent 生产部署")
    add_notes(
        slide,
        (
            "已证明的是问题真实、受众明确、规模可量化、底层旧版本有历史内网工程记录，"
            "以及当前版本能在本机完成真实 MPC 闭环。没有证明的是当前版本已部内复测或"
            "生产验收，也没有公开用户数、节省工时和人工修订率，所以这些指标不会编造。"
        ),
    )
    return slide


def appendix_quality(prs: Presentation, evidence: dict[str, Any], assets: dict[str, Path]):
    del evidence, assets
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    appendix_header(
        slide,
        "如何落地、扩展，并保证工程质量？",
        "直接结论：主链本地可运行、接口按区域解耦、质量检查和演示降级均已文档化。",
        "A5",
    )
    gates = [
        ("236", "镜像运行时测试", "green"),
        ("6", "宿主交付契约", "blue"),
        ("52", "接口兼容测试", "yellow"),
        ("3/3", "核心容器 healthy", "red"),
    ]
    for index, (value, label, accent) in enumerate(gates):
        metric_block(
            slide,
            Inches(0.7 + index * 3.13),
            Inches(1.92),
            Inches(2.72),
            value,
            label,
            "2026-07-30 决赛基线",
            accent,
        )

    columns = [
        (
            "本地落地",
            ["Gemma 4 + ADK 主链", "PostGIS + 优化引擎离线", "地图与 5 页图文 PDF"],
            "blue",
        ),
        (
            "按区域扩展",
            [
                "数据契约 / 坐标 / 编码检查",
                "区域适配器与工具接口解耦",
                "新区域独立硬约束校验",
            ],
            "green",
        ),
        (
            "质量与演示",
            ["Ruff / 编译 / 构建 / CI", "版本不兼容与缺产物安全停止", "实时 → 视频 → 封存产物"],
            "red",
        ),
    ]
    for index, (title, lines, accent) in enumerate(columns):
        x = Inches(0.8 + index * 4.13)
        rect(slide, x, Inches(3.42), Inches(3.5), Inches(2.45), "white", "line", radius=True)
        add_text(
            slide,
            x + Inches(0.25),
            Inches(3.7),
            Inches(2.98),
            Inches(0.26),
            title,
            14,
            accent,
            True,
        )
        add_text(
            slide,
            x + Inches(0.25),
            Inches(4.28),
            Inches(2.98),
            Inches(1.0),
            "\n".join(f"• {line}" for line in lines),
            11,
            "ink",
        )
    add_text(
        slide,
        Inches(1.05),
        Inches(6.25),
        Inches(11.0),
        Inches(0.3),
        "Google 栈口径：ADK 是已验证主运行时；AlphaEarth 仅 Tech Preview；OKF 仅知识交换 sidecar。",
        10,
        "navy",
        True,
        PP_ALIGN.CENTER,
    )
    footer(slide, "A5", "Google 技术的深度、必要性和证据比 Logo 数量更重要")
    add_notes(
        slide,
        (
            "当前主链可本地运行，并按数据契约、区域适配器和工具接口解耦扩展。工程证据"
            "分为 236 个镜像运行时测试、6 个宿主交付契约和 52 个接口兼容测试，不能直接"
            "混成一个总数。ADK 已进入主链；AlphaEarth 和 OKF 只作为未来扩展，不冒充现有能力。"
        ),
    )
    return slide


def build_deck(template: Path | None, output: Path) -> None:
    evidence = finals_evidence()
    assets = ensure_assets(evidence)
    if template and template.exists():
        reference = Presentation(str(template))
        slide_width = reference.slide_width
        slide_height = reference.slide_height
    else:
        slide_width = WIDE
        slide_height = HIGH
    # Reusing a template after deleting every slide can create duplicate part
    # names in python-pptx. Start clean and inherit only its 16:9 canvas.
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    builders = [
        slide_cover,
        slide_problem_solution,
        slide_nl2sql,
        slide_farmland_planning,
        slide_evidence,
        slide_close,
        appendix_architecture,
        appendix_reliability,
        appendix_claims,
        appendix_google,
        appendix_quality,
    ]
    for builder in builders:
        builder(prs, evidence, assets)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(json.dumps({"output": str(output), "slides": len(prs.slides)}, ensure_ascii=False))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--no-template", action="store_true")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    build_deck(None if args.no_template else args.template, args.output)


if __name__ == "__main__":
    main()
