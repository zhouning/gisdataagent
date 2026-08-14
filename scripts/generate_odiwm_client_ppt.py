#!/usr/bin/env python3
"""Generate the client-facing ODIWM exchange deck.

The deck is intentionally generated from local, versioned content so the
customer-facing statements and capability boundaries remain auditable.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = ROOT / "docs" / "reports"
ASSET_DIR = REPORT_DIR / "odiwm_customer_exchange_2026-08-14" / "assets"
DEFAULT_OUTPUT = REPORT_DIR / "odiwm_client_exchange_2026-08-14.pptx"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(22, 34, 54)
TEXT = RGBColor(47, 63, 82)
MUTED = RGBColor(92, 108, 126)
BG = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(204, 213, 222)
TEAL = RGBColor(19, 121, 108)
TEAL_PALE = RGBColor(231, 244, 241)
BLUE = RGBColor(47, 103, 174)
BLUE_PALE = RGBColor(232, 239, 248)
GREEN = RGBColor(100, 129, 48)
GREEN_PALE = RGBColor(239, 244, 232)
AMBER = RGBColor(190, 116, 13)
AMBER_PALE = RGBColor(255, 245, 224)
RED = RGBColor(171, 67, 62)
RED_PALE = RGBColor(249, 236, 235)
GRAY_PALE = RGBColor(239, 242, 245)

FONT_CN = "PingFang SC"
FONT_LATIN = "Aptos"


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
    font: str = FONT_CN,
    fit: bool = False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 1.08
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if fit:
        frame.fit_text(font_family=font, max_size=Pt(size))
    return box


def add_rich_text(slide, runs, x, y, w, h, *, size=18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    for text, color, bold in runs:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_CN
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LINE,
    radius: bool = True,
    width: float = 1.1,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, ValueError):
            pass
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=LINE, width=1.5, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_pill(slide, text, x, y, w, *, fill=TEAL_PALE, color=TEAL):
    add_box(slide, x, y, w, 0.34, fill=fill, line=fill)
    add_text(
        slide,
        text,
        x + 0.04,
        y + 0.02,
        w - 0.08,
        0.28,
        size=10.5,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_header(slide, number: str, title: str, subtitle: str | None = None):
    add_text(slide, number, 0.55, 0.34, 0.48, 0.32, size=12, color=TEAL, bold=True)
    add_text(slide, title, 1.02, 0.24, 11.7, 0.48, size=25, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, subtitle, 1.04, 0.73, 11.3, 0.34, size=11.5, color=MUTED)
    add_line(slide, 0.58, 1.08, 12.75, 1.08, color=LINE, width=0.8)


def add_footer(slide, page: int, label: str = "本体驱动的灌区世界模型 · 客户交流材料"):
    add_text(slide, label, 0.58, 7.12, 8.7, 0.18, size=8.5, color=MUTED)
    add_text(
        slide, f"{page:02d}", 12.28, 7.08, 0.48, 0.2, size=9, color=MUTED, align=PP_ALIGN.RIGHT
    )


def add_bullets(slide, items, x, y, w, h, *, size=16, color=TEXT, bullet_color=TEAL, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT_CN
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(gap)
        p.line_spacing = 1.13
        p.runs[0].font.color.rgb = color
    return box


def add_metric(slide, value, label, x, y, w, *, color=TEAL):
    add_text(slide, value, x, y, w, 0.48, size=26, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.48, w, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    add_box(slide, 0, 0, 13.333, 7.5, fill=NAVY, line=NAVY, radius=False)
    add_box(slide, 0.58, 0.72, 0.08, 4.75, fill=TEAL, line=TEAL, radius=False)
    add_pill(
        slide,
        "客户技术交流",
        0.9,
        0.78,
        1.35,
        fill=RGBColor(32, 58, 73),
        color=RGBColor(139, 218, 202),
    )
    add_text(slide, "本体驱动的灌区世界模型", 0.9, 1.42, 11.2, 0.7, size=36, color=WHITE, bold=True)
    add_text(
        slide,
        "从语义一致到可审查的条件推演",
        0.92,
        2.23,
        10.6,
        0.56,
        size=25,
        color=RGBColor(154, 220, 206),
        bold=True,
    )
    add_text(
        slide,
        "面向灌区状态认知、方案比较与调度辅助的系统设计",
        0.92,
        3.04,
        10.6,
        0.42,
        size=16,
        color=RGBColor(211, 220, 229),
    )
    add_line(slide, 0.92, 3.78, 11.98, 3.78, color=RGBColor(75, 91, 108), width=1)
    add_rich_text(
        slide,
        [
            ("本体", RGBColor(139, 218, 202), True),
            (" 负责确定对象、关系、状态与动作语义；  ", RGBColor(211, 220, 229), False),
            ("世界模型", RGBColor(144, 184, 234), True),
            (" 负责在约束下推演状态变化。", RGBColor(211, 220, 229), False),
        ],
        0.92,
        4.08,
        11.3,
        0.42,
        size=15,
    )
    add_text(
        slide,
        "技术交流材料  |  2026年8月14日",
        0.92,
        6.63,
        5.5,
        0.3,
        size=11,
        color=RGBColor(169, 182, 196),
    )
    add_text(
        slide,
        "Proposal only · 人工审查 · 不直接控制设备",
        7.1,
        6.58,
        5.15,
        0.36,
        size=11,
        color=RGBColor(169, 182, 196),
        align=PP_ALIGN.RIGHT,
    )
    return slide


def slide_position(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "01",
        "我们建议建设的不是另一个黑箱预测器",
        "目标是把业务语义、状态推演和责任链放到同一套可验证机制中",
    )
    add_box(slide, 0.72, 1.42, 3.56, 4.65, fill=TEAL_PALE, line=TEAL)
    add_pill(slide, "语义骨架", 0.98, 1.69, 1.15, fill=WHITE, color=TEAL)
    add_text(slide, "本体", 0.98, 2.18, 2.8, 0.48, size=28, color=TEAL, bold=True)
    add_text(slide, "确定系统在谈论什么", 0.98, 2.77, 2.8, 0.35, size=14.5, color=NAVY, bold=True)
    add_bullets(
        slide,
        [
            "稳定对象与权威拓扑",
            "State / Action / Constraint 合同",
            "单位、时间、空间、来源与版本",
            "自然语言落到确定的 Object + Link",
        ],
        0.98,
        3.25,
        2.9,
        2.28,
        size=13.5,
        color=TEXT,
    )

    add_box(slide, 4.88, 1.42, 3.56, 4.65, fill=BLUE_PALE, line=BLUE)
    add_pill(slide, "状态推演", 5.14, 1.69, 1.15, fill=WHITE, color=BLUE)
    add_text(slide, "世界模型", 5.14, 2.18, 2.8, 0.48, size=28, color=BLUE, bold=True)
    add_text(
        slide, "计算条件变化会带来什么", 5.14, 2.77, 2.9, 0.35, size=14.5, color=NAVY, bold=True
    )
    add_bullets(
        slide,
        [
            "冻结当前状态、外部驱动与候选动作",
            "守恒、容量、时延与专业模型",
            "比较未来状态、缺口和风险",
            "学习模型只在独立验证后作为增量",
        ],
        5.14,
        3.25,
        2.9,
        2.28,
        size=13.5,
        color=TEXT,
        bullet_color=BLUE,
    )

    add_box(slide, 9.04, 1.42, 3.56, 4.65, fill=AMBER_PALE, line=AMBER)
    add_pill(slide, "治理出口", 9.30, 1.69, 1.15, fill=WHITE, color=AMBER)
    add_text(slide, "Proposal", 9.30, 2.18, 2.8, 0.48, size=28, color=AMBER, bold=True)
    add_text(
        slide, "把结果交给有权人员审查", 9.30, 2.77, 2.9, 0.35, size=14.5, color=NAVY, bold=True
    )
    add_bullets(
        slide,
        [
            "候选动作、结果、风险与适用范围",
            "数据、模型、规则和运行版本",
            "可批准、驳回、修改与过期",
            "一期不直接触发闸门或泵站",
        ],
        9.30,
        3.25,
        2.9,
        2.28,
        size=13.5,
        color=TEXT,
        bullet_color=AMBER,
    )
    add_text(
        slide,
        "三者缺一不可：本体保证语义确定，世界模型负责条件推演，Proposal 保留组织责任。",
        1.33,
        6.35,
        10.7,
        0.42,
        size=15,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_business(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "02",
        "从四个高频问题开始，而不是从算法名称开始",
        "一期优先选择可获得数据、可建立基线、可由调度人员复核的场景",
    )
    cards = [
        ("状态看不全", "观测分散、缺测与时间错位，当前水量和供水缺口难统一判断", TEAL, TEAL_PALE),
        ("方案比较慢", "来水、设备或需求变化后，需要反复手工核算多个配水方案", BLUE, BLUE_PALE),
        ("结果难追溯", "方案依据、模型版本、规则和人工修改散落在不同记录中", AMBER, AMBER_PALE),
        ("经验难复用", "调度经验未沉淀为可查询的对象、规则、案例与证据", RED, RED_PALE),
    ]
    xs = [0.72, 3.86, 7.0, 10.14]
    for idx, ((title, body, color, pale), x) in enumerate(
        zip(cards, xs, strict=True), start=1
    ):
        add_box(slide, x, 1.52, 2.74, 2.55, fill=WHITE, line=LINE)
        add_box(slide, x + 0.18, 1.72, 0.42, 0.42, fill=color, line=color)
        add_text(
            slide,
            str(idx),
            x + 0.18,
            1.72,
            0.42,
            0.4,
            size=12,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, title, x + 0.72, 1.69, 1.78, 0.42, size=17, color=NAVY, bold=True)
        add_text(
            slide, body, x + 0.23, 2.35, 2.28, 1.25, size=12.5, color=TEXT, valign=MSO_ANCHOR.MIDDLE
        )
        add_box(slide, x + 0.18, 3.68, 2.38, 0.12, fill=pale, line=pale)
    add_text(slide, "一期目标", 0.74, 4.53, 1.1, 0.36, size=14, color=TEAL, bold=True)
    goals = [
        ("看得清", "统一对象、状态与质量"),
        ("算得明", "同条件比较有限候选"),
        ("审得过", "Proposal 全链路可追溯"),
        ("接得上", "衔接台账、报表与工单"),
    ]
    for i, (title, body) in enumerate(goals):
        x = 1.86 + i * 2.74
        add_box(slide, x, 4.42, 2.45, 1.12, fill=GRAY_PALE, line=LINE)
        add_text(slide, title, x + 0.14, 4.57, 0.76, 0.32, size=15, color=NAVY, bold=True)
        add_text(slide, body, x + 0.14, 4.96, 2.05, 0.30, size=11, color=MUTED)
    add_box(slide, 0.72, 5.92, 11.88, 0.73, fill=AMBER_PALE, line=AMBER)
    add_text(slide, "验收基线", 0.96, 6.10, 1.0, 0.30, size=13, color=AMBER, bold=True)
    add_text(
        slide,
        "与现行规则或人工方案做同任务、同时间窗比较；页面可用不等于模型有效。",
        2.02,
        6.08,
        9.95,
        0.33,
        size=13.5,
        color=NAVY,
        bold=True,
    )
    add_footer(slide, page)


def slide_ontology_value(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "03",
        "本体的核心价值：形成可计算、可治理的语义合同",
        "本体回答“是什么、连到哪里、允许做什么、证据来自哪里”，但不替代水动力计算",
    )
    items = [
        ("Object", "水源、渠段、闸门、量测点、田块、管理分区", TEAL, TEAL_PALE),
        ("Link", "upstreamOf、controlledBy、observedBy、supplies", BLUE, BLUE_PALE),
        ("State", "事实 / 估计 / 情景状态；时间、单位、质量与版本", GREEN, GREEN_PALE),
        ("Action", "开度、配水比例、启停、时段；对象、范围与权限", AMBER, AMBER_PALE),
        ("Constraint", "容量、水权、最低保障、时窗、设备与安全规则", RED, RED_PALE),
        ("Evidence", "数据来源、模型版本、适用域、审核与回执", NAVY, GRAY_PALE),
    ]
    for idx, (name, body, color, _pale) in enumerate(items):
        col = idx % 3
        row = idx // 3
        x = 0.72 + col * 4.03
        y = 1.42 + row * 2.12
        add_box(slide, x, y, 3.62, 1.76, fill=WHITE, line=color)
        add_box(slide, x, y, 0.13, 1.76, fill=color, line=color, radius=False)
        add_text(
            slide,
            name,
            x + 0.30,
            y + 0.22,
            1.25,
            0.38,
            size=19,
            color=color,
            bold=True,
            font=FONT_LATIN,
        )
        add_text(slide, body, x + 0.30, y + 0.77, 2.92, 0.65, size=12.3, color=TEXT)
    add_box(slide, 0.72, 5.80, 11.88, 0.78, fill=TEAL_PALE, line=TEAL)
    add_rich_text(
        slide,
        [
            ("关键控制：", TEAL, True),
            (
                "空间相交或名称相似不能自动认定为水力连通；权威拓扑必须核验、发布并版本冻结。",
                NAVY,
                True,
            ),
        ],
        0.98,
        6.03,
        11.3,
        0.34,
        size=13.5,
    )
    add_footer(slide, page)


def slide_ontology_mechanism(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "04",
        "本体如何进入运行链路：不是展示标签，而是约束每一步",
        "自然语言只负责表达意图；对象解析、数据访问和计算调用必须落到已注册合同",
    )
    stages = [
        ("问题", "未来24小时\n来水下降20%", NAVY, GRAY_PALE),
        ("Grounding", "范围 / 时间窗\n指标 / 对象类型", TEAL, TEAL_PALE),
        ("Object + Link", "R1 → C1\nC2/C3 → D1/D2", BLUE, BLUE_PALE),
        ("状态快照", "事实 + 估计\n质量 + 版本", GREEN, GREEN_PALE),
        ("Function", "守恒 / 时延\n专业模型 / 规划", AMBER, AMBER_PALE),
        ("Proposal", "动作 + 结果\n风险 + 证据", RED, RED_PALE),
    ]
    x0 = 0.45
    for idx, (title, body, color, pale) in enumerate(stages):
        x = x0 + idx * 2.12
        add_box(slide, x, 2.0, 1.72, 2.20, fill=pale, line=color)
        add_text(
            slide,
            title,
            x + 0.12,
            2.28,
            1.48,
            0.38,
            size=15.5,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            body,
            x + 0.13,
            2.96,
            1.46,
            0.72,
            size=11.5,
            color=TEXT,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if idx < len(stages) - 1:
            add_line(slide, x + 1.72, 3.10, x + 2.08, 3.10, color=MUTED, width=1.6, arrow=True)
    add_box(slide, 0.72, 4.76, 5.73, 1.28, fill=WHITE, line=TEAL)
    add_text(slide, "允许的路径", 0.98, 4.98, 1.15, 0.34, size=14, color=TEAL, bold=True)
    add_text(
        slide,
        "确定性检索对象和关系 → 读取有来源的状态 → 调用注册函数 → 生成待审核 Proposal",
        2.10,
        4.93,
        3.95,
        0.72,
        size=12.2,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_box(slide, 6.84, 4.76, 5.76, 1.28, fill=RED_PALE, line=RED)
    add_text(slide, "禁止的捷径", 7.10, 4.98, 1.15, 0.34, size=14, color=RED, bold=True)
    add_text(
        slide,
        "LLM 直接猜对象、拼 SQL、绕过函数合同或直接生成设备指令",
        8.25,
        4.93,
        3.95,
        0.72,
        size=12.2,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "LLM 是交互与编排入口，不是数据真实性、模型有效性或操作权限的来源。",
        1.45,
        6.40,
        10.4,
        0.34,
        size=14.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_world_model_value(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "05",
        "世界模型的核心价值：在明确条件下推演未来状态",
        "不预测像素；围绕灌区对象、网络关系和状态变量回答“如果这样做，会怎样”",
    )
    add_box(slide, 0.72, 1.43, 7.72, 2.28, fill=NAVY, line=NAVY)
    add_text(
        slide,
        "Sₜ₊₁ = F(Sₜ, Uₜ, Wₜ, θ  |  G, C)",
        1.06,
        1.82,
        7.05,
        0.72,
        size=29,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font=FONT_LATIN,
    )
    add_text(
        slide,
        "当前状态    候选动作    外部驱动    模型参数       拓扑关系    约束",
        1.08,
        2.73,
        7.0,
        0.35,
        size=11.8,
        color=RGBColor(188, 207, 224),
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "输出不是一句答案，而是未来状态轨迹、约束结果、不确定性与适用范围。",
        1.05,
        3.20,
        7.10,
        0.30,
        size=12.5,
        color=RGBColor(151, 220, 206),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_box(slide, 8.88, 1.43, 3.72, 2.28, fill=BLUE_PALE, line=BLUE)
    add_text(
        slide,
        "条件推演 ≠ 因果证明",
        9.17,
        1.78,
        3.08,
        0.42,
        size=19,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "动作收益需要历史回放、影子运行、对照设计或现场试验支持。模型只在已声明条件和适用域内给出结果。",
        9.22,
        2.39,
        2.98,
        0.93,
        size=12,
        color=TEXT,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    outputs = [
        ("到田水量", "各灌溉单元的供给与缺口"),
        ("传播过程", "时延、水位/流量或根区水分轨迹"),
        ("约束校核", "容量、最低保障、水权与动作合法性"),
        ("比较指标", "公平性、能耗、风险与现行方案差异"),
    ]
    for i, (title, body) in enumerate(outputs):
        x = 0.72 + i * 3.02
        add_box(slide, x, 4.25, 2.72, 1.55, fill=WHITE, line=LINE)
        add_text(
            slide,
            title,
            x + 0.20,
            4.52,
            2.30,
            0.34,
            size=15.5,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide, body, x + 0.23, 5.02, 2.24, 0.48, size=11.5, color=MUTED, align=PP_ALIGN.CENTER
        )
    add_box(slide, 0.72, 6.13, 11.88, 0.52, fill=AMBER_PALE, line=AMBER)
    add_text(
        slide,
        "模型有效性的前提：状态可识别、参数可校准、守恒可检查、独立时间窗可验证。",
        0.96,
        6.25,
        11.35,
        0.28,
        size=13.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_world_model_mechanism(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "06",
        "世界模型的实现机制：可信主线优先，学习模型按证据升级",
        "Geospatial Kernel 提供时空状态与算子执行基础，但与 JEPA 没有天然的架构等价关系",
    )
    layers = [
        ("状态估计", "多源观测对齐、缺测处理、事实/估计状态分离、不确定性", TEAL, TEAL_PALE),
        (
            "守恒与专业模型",
            "水量平衡、渠系时延；按数据条件选择 1D 水动力、渠池、根区水量平衡、FAO-56 等",
            BLUE,
            BLUE_PALE,
        ),
        (
            "规划与筛选",
            "规则、网络流、MILP / MPC 或适用搜索；先剔除硬约束违规方案",
            GREEN,
            GREEN_PALE,
        ),
        (
            "学习型增量",
            "降阶、图时空、神经算子或 JEPA 仅在同任务独立验证显示增益后进入",
            AMBER,
            AMBER_PALE,
        ),
    ]
    for idx, (title, body, color, pale) in enumerate(layers):
        y = 1.38 + idx * 1.22
        add_box(slide, 0.72, y, 8.22, 0.95, fill=pale, line=color)
        add_box(slide, 0.72, y, 1.62, 0.95, fill=color, line=color)
        add_text(
            slide,
            title,
            0.84,
            y + 0.24,
            1.38,
            0.36,
            size=14.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide, body, 2.58, y + 0.17, 6.04, 0.56, size=11.8, color=TEXT, valign=MSO_ANCHOR.MIDDLE
        )
    add_box(slide, 9.34, 1.38, 3.26, 4.61, fill=WHITE, line=NAVY)
    add_text(
        slide,
        "模型门禁",
        9.67,
        1.72,
        2.60,
        0.40,
        size=19,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    gates = [
        ("1", "可识别性"),
        ("2", "守恒与稳定性"),
        ("3", "独立窗口精度"),
        ("4", "速度与成本增益"),
        ("5", "适用域与失败保护"),
    ]
    for i, (n, label) in enumerate(gates):
        y = 2.38 + i * 0.62
        add_box(slide, 9.70, y, 0.38, 0.38, fill=NAVY, line=NAVY)
        add_text(
            slide,
            n,
            9.70,
            y,
            0.38,
            0.35,
            size=10.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, label, 10.27, y + 0.02, 1.65, 0.31, size=12.5, color=TEXT)
    add_box(slide, 0.72, 6.32, 11.88, 0.44, fill=GRAY_PALE, line=LINE)
    add_text(
        slide,
        "正确路线不是“Kernel 或 JEPA 二选一”，而是先建可信计算基线，"
        "再让学习模型在严格基准上证明增量。",
        0.96,
        6.40,
        11.35,
        0.26,
        size=13.1,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_closed_loop(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "07",
        "本体与世界模型如何形成一条可审查闭环",
        "本体确定计算边界；状态估计冻结输入；模型推演结果；Planner 比较候选；人工承担决策责任",
    )
    image = ASSET_DIR / "odiwm_architecture.png"
    slide.shapes.add_picture(
        str(image), Inches(0.74), Inches(1.30), width=Inches(8.08), height=Inches(4.73)
    )
    add_box(slide, 9.08, 1.30, 3.52, 4.73, fill=WHITE, line=LINE)
    add_text(slide, "边界判定", 9.39, 1.58, 2.85, 0.36, size=17, color=NAVY, bold=True)
    checks = [
        ("本体", "不计算水如何传播"),
        ("世界模型", "不决定组织权限"),
        ("Planner", "不把搜索结果称为全局最优"),
        ("LLM", "不直接接触原始库或设备"),
        ("人工", "审查适用性并保留最终决定"),
    ]
    for i, (title, body) in enumerate(checks):
        y = 2.11 + i * 0.68
        add_text(slide, title, 9.39, y, 0.78, 0.30, size=12.3, color=TEAL, bold=True)
        add_text(slide, body, 10.20, y, 2.05, 0.45, size=11.3, color=TEXT)
    add_box(slide, 0.74, 6.30, 11.86, 0.46, fill=TEAL_PALE, line=TEAL)
    add_text(
        slide,
        "闭环价值：每个结果都能回答“用了什么对象、状态、规则、模型和审核意见”。",
        0.98,
        6.39,
        11.3,
        0.28,
        size=13.5,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_proposal(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "08",
        "Proposal 把计算结果转化为可承担责任的候选方案",
        "动作与执行分离：模型产生 Proposal，有权人员审查，现有系统负责后续工单或控制流程",
    )
    add_box(slide, 0.72, 1.40, 7.70, 4.85, fill=WHITE, line=LINE)
    add_text(
        slide,
        "Proposal #IRR-20260814-003 · v2",
        1.02,
        1.70,
        4.6,
        0.38,
        size=18,
        color=NAVY,
        bold=True,
    )
    add_pill(slide, "待人工审查", 6.62, 1.70, 1.28, fill=AMBER_PALE, color=AMBER)
    rows = [
        ("问题与范围", "T+0~24h；R1-C1-C2/C3-D1/D2-F1~F4"),
        ("冻结输入", "ontology 0.1-demo · state snapshot v0.1 · scenario v0"),
        ("候选动作", "西支渠时段后移 6h；东/西分配比例调整为 45% / 55%"),
        ("结果摘要", "到田量、缺口、尾端保障、公平 CV、容量违规、水量账"),
        ("限制与风险", "合成参数；未校准；不代表真实灌区收益"),
        ("有效期", "仅限本次场景快照；状态或规则变化后必须重新运行"),
    ]
    for idx, (label, value) in enumerate(rows):
        y = 2.35 + idx * 0.57
        if idx:
            add_line(slide, 1.02, y - 0.08, 8.10, y - 0.08, color=LINE, width=0.6)
        add_text(slide, label, 1.02, y, 1.16, 0.33, size=11.5, color=MUTED, bold=True)
        add_text(slide, value, 2.30, y, 5.58, 0.42, size=11.3, color=TEXT)
    add_box(slide, 8.80, 1.40, 3.80, 1.24, fill=TEAL_PALE, line=TEAL)
    add_text(slide, "通过", 9.08, 1.67, 0.80, 0.34, size=16, color=TEAL, bold=True)
    add_text(slide, "记录审核意见，不触发设备", 9.08, 2.08, 3.12, 0.30, size=11.5, color=TEXT)
    add_box(slide, 8.80, 2.86, 3.80, 1.24, fill=AMBER_PALE, line=AMBER)
    add_text(slide, "退回", 9.08, 3.13, 0.80, 0.34, size=16, color=AMBER, bold=True)
    add_text(slide, "补充现场规则、状态或证据", 9.08, 3.54, 3.12, 0.30, size=11.5, color=TEXT)
    add_box(slide, 8.80, 4.32, 3.80, 1.24, fill=RED_PALE, line=RED)
    add_text(slide, "阻断", 9.08, 4.59, 0.80, 0.34, size=16, color=RED, bold=True)
    add_text(slide, "硬约束违规、输入失效或越权", 9.08, 5.00, 3.12, 0.30, size=11.5, color=TEXT)
    add_text(
        slide,
        "审查记录、修改版本和后续回执共同构成 Evidence。",
        1.35,
        6.52,
        10.7,
        0.32,
        size=14.2,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_demo(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "09",
        "今天的 Demo：真实服务架构，合成业务场景",
        "用于证明语义、推演、Proposal 与审核机制已经连通；不用于证明真实灌区预测精度",
    )
    image = ASSET_DIR / "odiwm_demo_flow.png"
    slide.shapes.add_picture(
        str(image), Inches(0.65), Inches(1.40), width=Inches(12.0), height=Inches(4.03)
    )
    add_box(slide, 0.72, 5.70, 3.63, 0.95, fill=TEAL_PALE, line=TEAL)
    add_text(slide, "后端权威运行", 0.98, 5.88, 1.30, 0.30, size=13, color=TEAL, bold=True)
    add_text(
        slide, "场景、公式、运行版本与审计在服务端", 2.25, 5.86, 1.80, 0.48, size=10.8, color=TEXT
    )
    add_box(slide, 4.55, 5.70, 3.63, 0.95, fill=BLUE_PALE, line=BLUE)
    add_text(slide, "API 驱动前端", 4.81, 5.88, 1.30, 0.30, size=13, color=BLUE, bold=True)
    add_text(
        slide, "参数、运行、回放与审核通过接口交互", 6.08, 5.86, 1.80, 0.48, size=10.8, color=TEXT
    )
    add_box(slide, 8.38, 5.70, 4.22, 0.95, fill=AMBER_PALE, line=AMBER)
    add_text(slide, "当前边界", 8.64, 5.88, 1.02, 0.30, size=13, color=AMBER, bold=True)
    add_text(
        slide, "合成参数 · 未校准 · Proposal only", 9.65, 5.86, 2.60, 0.48, size=10.8, color=TEXT
    )
    add_footer(slide, page)


def slide_system(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "10",
        "面向真实项目的系统结构：可替换数据与模型，不重写业务闭环",
        "一期把演示内核替换为客户数据、经校准模型和持久化审计；设备控制仍保持隔离",
    )
    columns = [
        (
            "业务交互",
            ["自然语言问题", "渠系与状态视图", "方案对比", "Proposal 审核"],
            NAVY,
            GRAY_PALE,
        ),
        (
            "应用服务",
            ["Ontology API", "Snapshot Service", "Scenario Run", "Proposal / Audit API"],
            TEAL,
            TEAL_PALE,
        ),
        (
            "计算运行时",
            ["状态估计", "守恒与专业模型", "Planner / Solver", "模型门禁与 Evals"],
            BLUE,
            BLUE_PALE,
        ),
        (
            "数据与治理",
            ["GIS 与工程台账", "时序观测与预报", "规则 / 水权 / 工单", "版本、来源与权限"],
            GREEN,
            GREEN_PALE,
        ),
    ]
    for idx, (title, items, color, pale) in enumerate(columns):
        x = 0.55 + idx * 3.17
        add_box(slide, x, 1.42, 2.76, 4.72, fill=WHITE, line=color)
        add_box(slide, x, 1.42, 2.76, 0.75, fill=color, line=color)
        add_text(
            slide,
            title,
            x + 0.16,
            1.62,
            2.44,
            0.32,
            size=15,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        for j, item in enumerate(items):
            y = 2.50 + j * 0.76
            add_box(slide, x + 0.22, y, 2.32, 0.52, fill=pale, line=pale)
            add_text(
                slide,
                item,
                x + 0.30,
                y + 0.10,
                2.16,
                0.30,
                size=11.6,
                color=TEXT,
                bold=j == 0,
                align=PP_ALIGN.CENTER,
            )
        if idx < 3:
            add_line(slide, x + 2.78, 3.76, x + 3.08, 3.76, color=MUTED, width=1.5, arrow=True)
    add_box(slide, 0.72, 6.38, 11.88, 0.42, fill=RED_PALE, line=RED)
    add_text(
        slide,
        "控制隔离：Proposal API 与设备 Action API 分离；一期仅提供审查输出，不下发真实控制指令。",
        0.95,
        6.45,
        11.40,
        0.26,
        size=12.8,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_data(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "11",
        "真实灌区落地的决定因素是数据合同与现场校准",
        "美国水库数据可用于研发验证，不作为中国灌区业务证据，也不在客户 Demo 中外发",
    )
    groups = [
        (
            "静态与拓扑",
            [
                "灌区边界、管理分区、渠系方向",
                "闸泵、分水口、量测点与控制关系",
                "田块、土壤、作物与供水关系",
            ],
            TEAL,
            TEAL_PALE,
        ),
        (
            "动态与外部驱动",
            [
                "流量、水位、闸位、工况与告警",
                "气象、降雨、遥感、土壤水分",
                "预报发布时间、修订时间与可用时间",
            ],
            BLUE,
            BLUE_PALE,
        ),
        (
            "规则与运行证据",
            [
                "水权、配额、轮灌、安全与例外规则",
                "计划、人工调整、审批与实际过流",
                "采纳、修改、驳回原因和事后结果",
            ],
            AMBER,
            AMBER_PALE,
        ),
    ]
    for idx, (title, items, color, pale) in enumerate(groups):
        x = 0.72 + idx * 4.02
        add_box(slide, x, 1.44, 3.62, 3.38, fill=WHITE, line=color)
        add_box(slide, x, 1.44, 3.62, 0.78, fill=pale, line=color)
        add_text(
            slide,
            title,
            x + 0.22,
            1.66,
            3.18,
            0.34,
            size=17,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_bullets(slide, items, x + 0.28, 2.52, 3.02, 1.82, size=12, color=TEXT, gap=8)
    add_text(slide, "每条状态至少携带", 0.74, 5.25, 1.62, 0.34, size=14, color=NAVY, bold=True)
    tags = ["时间", "单位", "来源", "质量", "空间范围", "版本", "可用时间"]
    for i, tag in enumerate(tags):
        add_pill(slide, tag, 2.50 + i * 1.25, 5.22, 0.98, fill=GRAY_PALE, color=NAVY)
    add_box(slide, 0.72, 5.92, 11.88, 0.80, fill=RED_PALE, line=RED)
    add_text(
        slide,
        "历史回放必须只使用当时实际可获得的数据",
        0.98,
        6.11,
        3.6,
        0.35,
        size=13.2,
        color=RED,
        bold=True,
    )
    add_text(
        slide,
        "禁止把事后修订的天气、观测或结果信息带入过去，否则会产生数据泄漏和虚假精度。",
        4.45,
        6.08,
        7.70,
        0.40,
        size=12.4,
        color=NAVY,
        bold=True,
    )
    add_footer(slide, page)


def slide_roadmap(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "12",
        "建议路线：先数字影子与条件推演，再决定是否扩大自动化",
        "每个阶段以进入条件为门禁；不以功能清单代替数据、模型和运行验证",
    )
    stages = [
        ("0", "现场确认", "场景、对象、指标与边界冻结", NAVY),
        ("1", "语义与数据", "Profile、权威拓扑、数据质量", TEAL),
        ("2", "数字影子", "历史/准实时状态与回放", BLUE),
        ("3", "条件推演", "校准模型、候选比较、Proposal", GREEN),
        ("4", "影子调度", "与现行调度并行，不影响生产", AMBER),
        ("5", "受控试点", "明确授权、人工复核、联合评审", RED),
    ]
    for idx, (num, title, body, color) in enumerate(stages):
        x = 0.43 + idx * 2.13
        add_box(slide, x, 2.02, 1.78, 2.52, fill=WHITE, line=color)
        add_box(slide, x + 0.56, 1.54, 0.66, 0.66, fill=color, line=color)
        add_text(
            slide,
            num,
            x + 0.56,
            1.66,
            0.66,
            0.38,
            size=15,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            title,
            x + 0.14,
            2.48,
            1.50,
            0.42,
            size=15,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            body,
            x + 0.18,
            3.22,
            1.42,
            0.78,
            size=10.8,
            color=TEXT,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if idx < 5:
            add_line(slide, x + 1.80, 3.27, x + 2.08, 3.27, color=MUTED, width=1.4, arrow=True)
    add_box(slide, 0.72, 5.10, 11.88, 1.20, fill=TEAL_PALE, line=TEAL)
    add_text(slide, "一期建议目标", 0.99, 5.38, 1.36, 0.34, size=15, color=TEAL, bold=True)
    add_text(
        slide,
        "完成 1–3 阶段：权威语义与拓扑、可回放的状态底座、"
        "经独立窗口验证的有限场景推演和可审查 Proposal。",
        2.38,
        5.30,
        9.72,
        0.55,
        size=13,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "跨灌区迁移、JEPA/神经算子、MCTS 与自动控制属于后续研究或受控试点，不作为一期默认承诺。",
        1.25,
        6.54,
        10.85,
        0.33,
        size=12.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_acceptance(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "13",
        "验收看证据，不看宣传口号",
        "核心标准是同任务比较、适用范围清楚、失败可发现、每一步可追溯",
    )
    left = [
        "对象、拓扑、单位、时间和来源映射正确",
        "缺测、重复、异常与过期状态可识别",
        "水量账残差、负存量、超容量可检测",
        "与现行规则或人工基线进行同任务比较",
    ]
    right = [
        "模型精度在独立时间窗评估",
        "Proposal 可追溯、可驳回、可修改、可过期",
        "计划、批准、执行与实际结果能够对账",
        "权限、人工确认和审计记录不可绕过",
    ]
    add_box(slide, 0.72, 1.45, 5.72, 4.65, fill=WHITE, line=TEAL)
    add_text(slide, "数据与模型", 1.02, 1.77, 4.95, 0.40, size=20, color=TEAL, bold=True)
    for i, item in enumerate(left):
        y = 2.43 + i * 0.80
        add_box(slide, 1.02, y, 0.38, 0.38, fill=TEAL, line=TEAL)
        add_text(
            slide,
            "✓",
            1.02,
            y - 0.01,
            0.38,
            0.37,
            size=13,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, item, 1.58, y, 4.30, 0.50, size=12.3, color=TEXT)
    add_box(slide, 6.88, 1.45, 5.72, 4.65, fill=WHITE, line=BLUE)
    add_text(slide, "运行与治理", 7.18, 1.77, 4.95, 0.40, size=20, color=BLUE, bold=True)
    for i, item in enumerate(right):
        y = 2.43 + i * 0.80
        add_box(slide, 7.18, y, 0.38, 0.38, fill=BLUE, line=BLUE)
        add_text(
            slide,
            "✓",
            7.18,
            y - 0.01,
            0.38,
            0.37,
            size=13,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, item, 7.74, y, 4.30, 0.50, size=12.3, color=TEXT)
    add_box(slide, 0.72, 6.36, 11.88, 0.42, fill=AMBER_PALE, line=AMBER)
    add_text(
        slide,
        "“毫秒级、1000×、全局最优、因果收益”只有在任务口径、硬件、基线和独立试验明确后才可成为指标。",
        0.94,
        6.43,
        11.45,
        0.27,
        size=12.4,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def slide_questions(prs, page):
    slide = blank_slide(prs)
    add_header(
        slide,
        "14",
        "建议今天共同确认七个问题",
        "这些答案决定一期范围、数据准备、模型选择和可验收结果",
    )
    questions = [
        "首期最有价值的场景：短期缺口、轮灌调整、故障应对，还是其他问题？",
        "可提供哪些渠系、量测、田块、作物和调度日志？时间跨度和质量如何？",
        "哪些对象与关系已有权威来源，哪些需要现场核验？",
        "当前方案可行性由谁判断，使用哪些指标、规则和经验？",
        "哪些结果可以作为建议，哪些必须由特定角色审批？",
        "是否允许先做历史回放和不影响生产的影子运行？",
        "数据共享、脱敏、部署、审计留存和后续试点有哪些约束？",
    ]
    for idx, question in enumerate(questions):
        col = 0 if idx < 4 else 1
        row = idx if idx < 4 else idx - 4
        x = 0.72 if col == 0 else 6.84
        y = 1.38 + row * 1.26
        add_box(slide, x, y, 5.76, 0.96, fill=WHITE, line=LINE)
        add_box(
            slide,
            x + 0.18,
            y + 0.22,
            0.48,
            0.48,
            fill=TEAL if col == 0 else BLUE,
            line=TEAL if col == 0 else BLUE,
        )
        add_text(
            slide,
            str(idx + 1),
            x + 0.18,
            y + 0.24,
            0.48,
            0.39,
            size=11.5,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            question,
            x + 0.83,
            y + 0.15,
            4.58,
            0.64,
            size=11.8,
            color=NAVY,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    add_box(slide, 6.84, 5.16, 5.76, 1.00, fill=TEAL_PALE, line=TEAL)
    add_text(slide, "建议下一步", 7.12, 5.42, 1.18, 0.35, size=14.5, color=TEAL, bold=True)
    add_text(
        slide,
        "选择 1 个场景 + 1 个典型渠系 + 1 段可回放历史窗口，完成联合数据审查。",
        8.32,
        5.30,
        3.91,
        0.55,
        size=12,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "目标不是让模型替代调度人员，而是让每一次方案比较更清楚、更快、更可追溯。",
        1.23,
        6.55,
        10.87,
        0.34,
        size=14.3,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page)


def closing_slide(prs, page):
    slide = blank_slide(prs)
    add_box(slide, 0, 0, 13.333, 7.5, fill=NAVY, line=NAVY, radius=False)
    add_pill(
        slide,
        "讨论起点",
        0.92,
        0.86,
        1.12,
        fill=RGBColor(32, 58, 73),
        color=RGBColor(139, 218, 202),
    )
    add_text(slide, "先把语义和证据做实，", 0.92, 1.67, 10.8, 0.67, size=32, color=WHITE, bold=True)
    add_text(
        slide,
        "再让模型在边界内变快、变准。",
        0.92,
        2.42,
        11.25,
        0.67,
        size=32,
        color=RGBColor(154, 220, 206),
        bold=True,
    )
    add_line(slide, 0.92, 3.53, 11.95, 3.53, color=RGBColor(75, 91, 108), width=1)
    add_text(
        slide, "本体", 0.94, 4.02, 1.3, 0.40, size=18, color=RGBColor(139, 218, 202), bold=True
    )
    add_text(
        slide,
        "确定对象、关系、状态、动作与证据",
        2.12,
        4.03,
        4.2,
        0.35,
        size=14,
        color=RGBColor(211, 220, 229),
    )
    add_text(
        slide, "世界模型", 0.94, 4.67, 1.3, 0.40, size=18, color=RGBColor(144, 184, 234), bold=True
    )
    add_text(
        slide,
        "在守恒、专业模型和约束下推演条件变化",
        2.12,
        4.68,
        5.1,
        0.35,
        size=14,
        color=RGBColor(211, 220, 229),
    )
    add_text(
        slide, "治理闭环", 0.94, 5.32, 1.3, 0.40, size=18, color=RGBColor(234, 186, 116), bold=True
    )
    add_text(
        slide,
        "Proposal、人工审查、执行隔离与全链路审计",
        2.12,
        5.33,
        5.1,
        0.35,
        size=14,
        color=RGBColor(211, 220, 229),
    )
    add_text(
        slide,
        "Q & A",
        9.88,
        5.67,
        2.08,
        0.72,
        size=31,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.RIGHT,
        font=FONT_LATIN,
    )
    add_text(
        slide,
        f"{page:02d}",
        12.28,
        7.08,
        0.48,
        0.2,
        size=9,
        color=RGBColor(169, 182, 196),
        align=PP_ALIGN.RIGHT,
    )


def set_metadata(prs):
    props = prs.core_properties
    props.title = "本体驱动的灌区世界模型：从语义一致到可审查的条件推演"
    props.subject = "ODIWM 客户技术交流方案"
    props.author = "技术交流材料"
    props.keywords = "灌区, 本体, 世界模型, 条件推演, Proposal, 人工审查"
    props.comments = "客户交流版；合成场景不代表真实灌区模型能力。"


def build(output: Path):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    set_metadata(prs)
    title_slide(prs)
    slide_position(prs, 2)
    slide_business(prs, 3)
    slide_ontology_value(prs, 4)
    slide_ontology_mechanism(prs, 5)
    slide_world_model_value(prs, 6)
    slide_world_model_mechanism(prs, 7)
    slide_closed_loop(prs, 8)
    slide_proposal(prs, 9)
    slide_demo(prs, 10)
    slide_system(prs, 11)
    slide_data(prs, 12)
    slide_roadmap(prs, 13)
    slide_acceptance(prs, 14)
    slide_questions(prs, 15)
    closing_slide(prs, 16)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return len(prs.slides)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    slide_count = build(args.output.resolve())
    print(f"generated {slide_count} slides: {args.output.resolve()}")


if __name__ == "__main__":
    main()
